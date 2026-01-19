import io
import base64
import uuid
import re
import os
from matplotlib import image
import requests
import dotenv
import psycopg2
import fitz  # PyMuPDF
import time
import mimetypes # For guessing mime types
import numpy as np

from PIL import Image
# from sentence_transformers import SentenceTransformer
# from transformers import BitsAndBytesConfig, AutoModel, AutoProcessor
# import torch
# from colpali_engine.models import ColIdefics3, ColIdefics3Processor

# New imports for Docling and concurrent processing
# from docling_core.types.doc import ImageRefMode, PictureItem, TableItem
# from docling.datamodel.base_models import InputFormat
# from docling.document_converter import DocumentConverter, PdfFormatOption
# from docling.datamodel.base_models import DocumentStream
# from docling_core.types.doc import DocItemLabel
# import concurrent.futures
# from docling.pipeline.vlm_pipeline import VlmPipeline
# from docling.datamodel.pipeline_options import (
#     VlmPipelineOptions,
#     PdfPipelineOptions,
#     TesseractCliOcrOptions,
#     TesseractOcrOptions,
#     RapidOcrOptions,
# )
# from modelscope import snapshot_download
# from docling.datamodel.accelerator_options import AcceleratorDevice, AcceleratorOptions
# from docling.datamodel import vlm_model_specs
# from docling.datamodel.pipeline_options_vlm_model import ApiVlmOptions, ResponseFormat
# from docling.pipeline.vlm_pipeline import VlmPipeline

import base64
import io
# from unstructured.partition.auto import partition
# from unstructured.documents.elements import Image as UnstructuredImage

from minio import Minio
from minio.error import S3Error
from typing import List, Optional, Dict, Any, Union

import openai
import httpx

import os
from io import BytesIO
from docx import Document
from pptx import Presentation
import openpyxl
import xlrd
from striprtf.striprtf import rtf_to_text
from odf.opendocument import load
from odf import text


# from vllm import LLM, SamplingParams
# from vllm.config import PoolerConfig
# from vllm.inputs.data import TextPrompt

dotenv.load_dotenv()

LOCAL = os.getenv("LOCAL", True)
IFXGPT = os.getenv("IFXGPT", True)
API_OLLAMA = os.getenv("API_OLLAMA", "http://127.0.0.1:11434/api/generate")

# IFX GPT API
cert_path = os.getenv('IFXGPT_CERT_PATH')   
window_user = os.getenv('WINDOWS_USER')
window_password = os.getenv('WINDOWS_PASSWORD') 
token = os.getenv("TOKEN")
basic_auth = False

LOCAL = True if LOCAL == "True" else False
IFXGPT = True if IFXGPT == "True" else False

def generate_base64_string(username, password):   
    """   
    To Encode username & password strings into base64   
    """   
   
    sample_string = username + ":" + password   
    sample_string_bytes = sample_string.encode("ascii")   
   
    base64_bytes = base64.b64encode(sample_string_bytes)   
    base64_string = base64_bytes.decode("ascii")   
    return base64_string   
   
if basic_auth:   
    token = generate_base64_string(window_user, window_password)     
    headers = {      
    'Authorization': f"Basic {token}"      
    }   
else:    
    token = '<generated-token-from-url>'    
    headers = {     
        'Authorization': f"Bearer {token}",      }      
   
client = openai.OpenAI(     
            api_key=token,     
            base_url='https://gpt4ifx.icp.infineon.com',      
            default_headers=headers,     
            http_client = httpx.Client(verify=cert_path)   
            )    

# Function to calculate parameter memory
def get_model_memory(model):
    total_params = sum(p.numel() for p in model.parameters())  # Total number of parameters
    total_memory_bytes = sum(p.numel() * p.element_size() for p in model.parameters())  # In bytes
    total_memory_mb = total_memory_bytes / (1024 ** 2)  # Convert to MB
    total_memory_gb = total_memory_mb / 1024  # Convert to GB
    
    print(f"Total Parameters: {total_params:,}")
    print(f"Total Memory: {total_memory_gb:.2f} GB ({total_memory_mb:.2f} MB)")
    return total_memory_gb

def clear_gpu():
    import torch
    import gc
    if torch.cuda.is_available():
        torch.cuda.empty_cache()
        torch.cuda.ipc_collect()
        print("Cleared GPU memory.")
    gc.collect()

# ==============================================================================
#  DATABASE & MINIO CLIENT SETUP
# ==============================================================================
# clear_gpu()
# device = "cuda" if torch.cuda.is_available() else "cpu"

# Initialize model variable (will be set below if LOCAL=True)
model = None

# Quantization model
if LOCAL:
    print("🚀 Loading Jina embedding model (LOCAL mode)...")
    # 1. Define the 4-bit configuration
    bnb_config = BitsAndBytesConfig(
        load_in_4bit=True,
        bnb_4bit_quant_type="nf4",      # Normalized Float 4 (best for accuracy)
        bnb_4bit_compute_dtype=torch.bfloat16, # Use bfloat16 for computation (requires Ampere+ GPU), otherwise use float16
        bnb_4bit_use_double_quant=True  # Saves a bit more memory
    )

    # 2. Load the model with the config
    try:
        model = SentenceTransformer(
            "jinaai/jina-embeddings-v4",
            trust_remote_code=True,
            model_kwargs={
                "default_task": "retrieval",
                "quantization_config": bnb_config,
                "device_map": "auto"  # REQUIRED: Lets accelerate handle GPU placement
            },
        )
        clear_gpu()
        print("✅ Jina model loaded successfully")
        print(model)
        uses_mem = get_model_memory(model)
    except Exception as e:
        print(f"❌ Failed to load Jina model: {e}")
        model = None
else:
    print("📡 Running in REMOTE mode (will use IFXGPT/API for embeddings)")

# --- NEW: MinIO Client Initialization ---
minio_client = Minio(
    os.getenv("MINIO_ENDPOINT", "127.0.0.1:9010") + ":" + os.getenv("MINIO_PORT", "9010"),
    access_key=os.getenv("MINIO_ACCESS_KEY", "minioadmin"),
    secret_key=os.getenv("MINIO_SECRET_KEY", "minioadmin"),
    secure=os.getenv("MINIO_USE_SSL", "false").lower() == 'true'
)
minio_bucket_name = os.getenv("MINIO_BUCKET", "user-files")

# --- NEW: Database Connection Helper ---
def get_db_connection():
    """Helper function to get a new database connection."""
    try:
        conn = psycopg2.connect(
            dbname=os.getenv("PGDATABASE", "ai_agent"),
            user=os.getenv("PGUSER", "athip"),
            password=os.getenv("PGPASSWORD", "123456"),
            host=os.getenv("PGHOST", "localhost"),
            port=os.getenv("PGPORT", "5432"),
        )
        return conn
    except psycopg2.DatabaseError as e:
        print(f"❌ Failed to connect to database: {e}")
        return None

# ==============================================================================
#  NEW: MINIO & FILE PROCESSING HELPERS
# ==============================================================================

def upload_file_to_minio_and_db(user_id: int, chat_history_id: int, file_name: str, file_bytes: bytes, mime_type: str = None) -> tuple[Optional[int], Optional[str]]:
    """
    Uploads a file to MinIO and creates a record in the 'uploaded_files' table.
    Mirrors the logic from database.js.

    Returns:
        (uploaded_file_id, object_name) or (None, None) on failure.
    """
    if not mime_type:
        mime_type, _ = mimetypes.guess_type(file_name)
        if not mime_type:
            mime_type = 'application/octet-stream' # Default
            
    file_size = len(file_bytes)
    object_name = f"user_{user_id}/chat_{chat_history_id}/{int(time.time())}-{file_name}"
    
    print(f"📦 DEBUG upload_file_to_minio_and_db:")
    print(f"  - user_id: {user_id}")
    print(f"  - chat_history_id: {chat_history_id}")
    print(f"  - file_name: {file_name}")
    print(f"  - file_size: {file_size} bytes")
    print(f"  - mime_type: {mime_type}")
    print(f"  - object_name: {object_name}")
    
    conn = None
    try:
        # 1. Upload to MinIO
        print("  - Attempting MinIO upload...")
        file_stream = io.BytesIO(file_bytes)
        minio_client.put_object(
            minio_bucket_name,
            object_name,
            file_stream,
            file_size,
            content_type=mime_type
        )
        print(f"✅ MinIO: File '{object_name}' uploaded successfully.")

        # 2. Insert record into PostgreSQL
        print("  - Attempting DB insert...")
        conn = get_db_connection()
        if not conn:
            raise Exception("Could not connect to database")
            
        cur = conn.cursor()
        query = """
            INSERT INTO uploaded_files (user_id, chat_history_id, file_name, object_name, mime_type, file_size_bytes)
            VALUES (%s, %s, %s, %s, %s, %s)
            RETURNING id;
        """
        values = (user_id, chat_history_id, file_name, object_name, mime_type, file_size)
        cur.execute(query, values)
        uploaded_file_id = cur.fetchone()[0]
        
        conn.commit()
        cur.close()
        
        print(f"✅ DB: Saved file record with ID: {uploaded_file_id}")
        return uploaded_file_id, object_name

    except Exception as e:
        print(f"❌ Error during file upload process: {e}")
        import traceback
        traceback.print_exc()
        # Attempt to clean up MinIO object if DB insert fails or MinIO fails
        try:
            minio_client.remove_object(minio_bucket_name, object_name)
            print(f"🧹 Cleaned up MinIO object '{object_name}' after error.")
        except S3Error as cleanup_error:
            print(f"⚠️ Failed to clean up MinIO object '{object_name}': {cleanup_error}")
        if conn:
            conn.rollback()
        return None, None
    finally:
        if conn:
            conn.close()


def get_file_from_minio(object_name: str) -> Optional[bytes]:
    """Retrieves a file's content from MinIO as bytes."""
    try:
        response = minio_client.get_object(minio_bucket_name, object_name)
        file_bytes = response.read()
        return file_bytes
    except S3Error as e:
        print(f"❌ Error getting file '{object_name}' from MinIO: {e}")
        return None
    finally:
        if 'response' in locals() and response:
            response.close()
            response.release_conn()

def convert_pdf_page_to_image(pdf_bytes: bytes, page_number_0_indexed: int, dpi: int = 100) -> Optional[bytes]:
    """
    Extracts a single page from a PDF as a high-quality PNG image.
    
    Args:
        pdf_bytes: The byte content of the entire PDF.
        page_number_0_indexed: The page to extract (0 for page 1, 1 for page 2, etc.)
        dpi: The resolution in dots per inch for the output image.
    """
    try:
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
        if page_number_0_indexed >= len(pdf_document):
            print(f"Error: Page {page_number_0_indexed} out of bounds.")
            pdf_document.close()
            return None
            
        page = pdf_document.load_page(page_number_0_indexed)
        
        # Render page to a pixmap (image)
        # Use high DPI (e.g., 200) for good quality
        pix = page.get_pixmap(dpi=dpi) 
        
        # Convert pixmap to PNG bytes
        img_bytes = pix.tobytes("png")
        pdf_document.close()
        return img_bytes
        
    except Exception as e:
        print(f"❌ Error converting PDF page {page_number_0_indexed} to image: {e}")
        return None


# ==============================================================================
#  LEGACY: DOCLING / UNSTRUCTURED (No changes from your code)
# ==============================================================================

def extract_images_from_pdf(pdf_bytes: bytes) -> list:
    """
    Extracts all images from PDF bytes using PyMuPDF and returns them as a list of bytes.
    """
    images = []
    try:
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
        print(f"📄 PDF has {pdf_document.page_count} pages. Extracting images...")
        for page_number in range(len(pdf_document)):
            page = pdf_document.load_page(page_number)
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                images.append(image_bytes)
        print(f"🖼️ Successfully extracted {len(images)} images from the PDF.")
    except Exception as e:
        print(f"❌ Error extracting images with PyMuPDF: {e}")
    return images


def extract_images_from_pdf_with_docling(pdf_bytes: bytes) -> list:
    """
    Extracts images from a PDF using the Docling pipeline.
    This method is more accurate as it identifies semantic PictureItems
    instead of just raw image objects like PyMuPDF.
    """
    images = []
    doc_stream = DocumentStream(name="temp_doc.pdf", stream=io.BytesIO(pdf_bytes))
    
    print("🧠 Using Docling model to identify and extract images...")
    
    try:
        # --- Configure a standard OCR pipeline for layout analysis ---
        # This setup is efficient for finding image blocks without running a full VLM.
        # print("Downloading RapidOCR models for layout analysis...")
        # download_path = snapshot_download(repo_id="RapidAI/RapidOCR")
        # det_model_path = os.path.join(download_path, "onnx", "PP-OCRv5", "det", "ch_PP-OCRv5_server_det.onnx")
        # rec_model_path = os.path.join(download_path, "onnx", "PP-OCRv5", "rec", "ch_PP-OCRv5_rec_server_infer.onnx")
        # cls_model_path = os.path.join(download_path, "onnx", "PP-OCRv4", "cls", "ch_ppocr_mobile_v2.0_cls_infer.onnx")
        
        # ocr_options = RapidOcrOptions(
        #     det_model_path=det_model_path, rec_model_path=rec_model_path, cls_model_path=cls_model_path
        # )

        # pipeline_options = PdfPipelineOptions()
        # pipeline_options.do_ocr = True # OCR helps in distinguishing text from pure image areas
        # pipeline_options.ocr_options = ocr_options


        pipeline_options = VlmPipelineOptions(
        vlm_options=vlm_model_specs.SMOLDOCLING_TRANSFORMERS,  # <-- change the model here
            # accelerator_options = AcceleratorOptions(
            #     num_threads=8, device=AcceleratorDevice.CPU
            # )
        )
        
        converter = DocumentConverter(
            format_options={InputFormat.PDF: PdfFormatOption(pipeline_cls=VlmPipeline, pipeline_options=pipeline_options)}
        )

        # converter = DocumentConverter()
        
        # --- Convert the document and iterate through items ---
        result = converter.convert(doc_stream)
        doc = result.document

        for item, _level in doc.iterate_items():
            # Check if the item is a PictureItem
            if isinstance(item, PictureItem):
                # Get the image resource (as a PIL Image)
                image_resource = item.get_image(doc)
                
                # Convert the PIL Image to bytes
                buf = io.BytesIO()
                image_resource.save(buf, format="PNG") # Save as PNG for consistency
                image_bytes = buf.getvalue()
                images.append(image_bytes)
        
        print(f"🖼️ Successfully extracted {len(images)} images using the Docling model.")
        
    except Exception as e:
        print(f"❌ Error extracting images with Docling: {e}")
        
    return images


def extract_images_from_doc_with_unstructured(doc_bytes: bytes, filename: str) -> list:
    """
    Extracts images from a document using the 'unstructured' library.
    """
    images = []
    print(f"🧠 Using 'unstructured' library with hi_res strategy to extract images from '{filename}'...")

    try:
        # 🛠️ Configure the partitioner specifically for image extraction
        elements = partition(
            file=io.BytesIO(doc_bytes),
            file_filename=filename,
            # extract_images_in_pdf=True, 
            extract_image_block_types=['Image'],
            extract_image_block_to_payload=True,
            # extract_image_block_output_dir="./temp_images",
            infer_table_structure=True,
            languages=['tha', 'eng'],
            strategy="hi_res",
            hi_res_model_name="yolox"
        )

        # Iterate through the elements found in the document
        for element in elements:
            # Check if the element is an image and has the embedded base64 data
            if isinstance(element, UnstructuredImage) and element.metadata.image_base64:
                # Decode the base64 string back into raw image bytes
                image_bytes = base64.b64decode(element.metadata.image_base64)
                images.append(image_bytes)

        print(f"🖼️ Successfully extracted {len(images)} images using 'unstructured'.")

    except Exception as e:
        print(f"❌ Error extracting images with 'unstructured': {e}")

    return images


# ==============================================================================
#  LEGACY: VLM & CONTENT EXTRACTION (No changes from your code)
# ==============================================================================

def _create_deepinfra_vlm_options(model: str, prompt: str, api_key: str):
    """
    Helper function to create ApiVlmOptions specifically for DeepInfra's OpenAI-compatible endpoint.
    """
    if not api_key:
        raise ValueError("DeepInfra API key is required.")

    headers = {"Authorization": f"Bearer {api_key}"}
    print(f'prompt: {prompt[:1000]}')

    # Configuration for the remote VLM API call
    options = ApiVlmOptions(
        url="https://api.deepinfra.com/v1/openai/chat/completions",
        params=dict(
            model=model,
            max_tokens=4096,  # Maximum tokens for the response
        ),
        headers=headers,
        prompt=prompt,
        timeout=180,  # Increased timeout for potentially large documents
        scale=1.0,  # Scale down images to reduce payload size
        temperature=0.0,  # Lower temperature for more deterministic, factual output
        # MARKDOWN is a direct and useful output. DOCTAGS is another option for a more structured format.
        response_format=ResponseFormat.MARKDOWN,
        concurrency=10,  # Allow multiple concurrent requests
    )
    return options

def _create_openrouter_vlm_options(model: str, prompt: str, api_key: str):
    """
    Helper function to create ApiVlmOptions specifically for OpenRouter's API.
    """
    if not api_key:
        raise ValueError("OpenRouter API key is required.")

    headers = {"Authorization": f"Bearer {api_key}"}
    print(f'prompt: {prompt[:1000]}')

    # Configuration for the remote VLM API call
    options = ApiVlmOptions(
        url="https://openrouter.ai/api/v1/chat/completions",
        params=dict(
            model=model,
            max_tokens=4096*10,  # Maximum tokens for the response
        ),
        headers=headers,
        prompt=prompt,
        timeout=180,  # Increased timeout for potentially large documents
        scale=2.0,  # Scale down images to reduce payload size
        temperature=0.0,  # Lower temperature for more deterministic, factual output
        # MARKDOWN is a direct and useful output. DOCTAGS is another option for a more structured format.
        response_format=ResponseFormat.MARKDOWN,
        concurrency=15,  # Allow multiple concurrent requests
    )
    return options


def generate_vlm_pipeline_options(mode: str = 'local', **kwargs):
    """
    Generates pipeline_options for the Docling VLM pipeline based on the specified mode.
    """
    if mode == 'local':
        model_name = kwargs.get('model_name', 'SMOLDOCLING_TRANSFORMERS')
        try:
            model_spec = getattr(vlm_model_specs, model_name)
        except AttributeError:
            raise ValueError(f"Model '{model_name}' not found in docling.datamodel.vlm_model_specs.")

        print(f"✅ Configuring VLM pipeline for LOCAL mode with model: {model_name}")
        pipeline_options = VlmPipelineOptions(
            vlm_options=model_spec,
        )
        return pipeline_options

    elif mode == 'remote':
        model_name = kwargs.get('model_name', 'deepseek-ai/DeepSeek-OCR')
        prompt = kwargs.get('prompt', 'Convert this document page to clean, well-structured markdown. Extract all text, format tables, and describe any images or charts in detail.')
        api_key = kwargs.get('api_key') or os.getenv("OPENROUTER_API_KEY") #os.getenv("DEEPINFRA_API_KEY")
        # api_key = kwargs.get('api_key') or os.getenv("DEEPINFRA_API_KEY")
        if model_name.startswith('deepseek-ai/DeepSeek-OCR'):
            prompt = f"<|begin_of_text|><|start_header_id|>user<|end_header_id|>\n\n{prompt}<|eot_id|><|start_header_id|>assistant<|end_header_id|>\n\n"
        elif model_name.startswith('allenai/olmOCR-2-7B-1025'):
            prompt = ""

        if not api_key:
            raise ValueError("OPENROUTER_API_KEY not found. Please provide it as a kwarg or set the environment variable.")

        print(f"✅ Configuring VLM pipeline for REMOTE mode with OpenRouter model: {model_name}")
        api_vlm_options = _create_openrouter_vlm_options(
            model=model_name,
            prompt=prompt,
            api_key=api_key
        )

        # print(f"✅ Configuring VLM pipeline for REMOTE mode with OpenRouter model: {model_name}")
        # api_vlm_options = _create_deepinfra_vlm_options(
        #     model=model_name,
        #     prompt=prompt,
        #     api_key=api_key
        # )        

        pipeline_options = VlmPipelineOptions(
            enable_remote_services=True,  # This is required to allow API calls
            vlm_options=api_vlm_options,
            do_picture_description=True,  # Ensure images are described
            images_scale=2
        )
        return pipeline_options

    else:
        raise ValueError(f"Invalid mode '{mode}' specified. Choose 'local' or 'remote'.")
    

# def extract_and_process_content(file_storage, option: str = 'describe', pipeline_mode: str = 'ocr', with_images: bool = True) -> str:
#     """
#     Extracts content from a file, gets image descriptions, and replaces
#     image placeholders in the markdown output with those descriptions.
    
#     NOTE: This is part of the 'legacy' flow.
#     """
#     if option not in ['describe', 'summarize', 'describe_images_only']:
#         return "Error: Invalid option provided."

#     # This function expects a file_storage object, so we wrap bytes in BytesIO
#     if isinstance(file_storage, bytes):
#         file_bytes = file_storage
#         file_storage = io.BytesIO(file_bytes)
#         file_storage.filename = "temp_file.pdf" # Mock filename
#     else:
#         file_storage.seek(0)
#         file_bytes = file_storage.read()
#         file_storage.seek(0) # Reset pointer

#     # --- Step 1: Extract images first using PyMuPDF ---
#     extracted_images_bytes = []
#     if file_storage.filename.lower().endswith('.pdf') and with_images:
#         #  extracted_images_bytes = extract_images_from_pdf(file_bytes)
#          extracted_images_bytes = extract_images_from_pdf_with_docling(file_bytes)
#         #  extracted_images_bytes = extract_images_from_doc_with_unstructured(
#         #                                                                     doc_bytes=file_bytes,
#         #                                                                     filename=file_storage.filename
#         #                                                                 )
#          extracted_images_bytes = list(reversed(extracted_images_bytes))

#     # Handle the case where only image descriptions are requested
#     if option == 'describe_images_only':
#         if not extracted_images_bytes:
#             return "No images found in the document."
#         with concurrent.futures.ThreadPoolExecutor() as executor:
#             future_to_desc = {executor.submit(image_to_describe_from_base64, img_bytes): i for i, img_bytes in enumerate(extracted_images_bytes)}
#             descriptions = [future.result() for future in concurrent.futures.as_completed(future_to_desc)]
#         return "\n\n--- IMAGE DESCRIPTION ---\n".join(descriptions)

#     # --- Step 2: Run the Docling pipeline to get text ---
#     doc_stream = DocumentStream(name=file_storage.filename, stream=io.BytesIO(file_bytes))
#     print(f"📄 Converting document with Docling using '{pipeline_mode}' pipeline...")
    
#     try:
#         # (The setup logic for the 'converter' object remains unchanged)
#         if pipeline_mode == 'ocr':
#            # --- ORIGINAL OCR-BASED PIPELINE SETUP ---
#             # Download RapidOCR models from Hugging Face
#             print("Downloading RapidOCR models")
#             download_path = snapshot_download(repo_id="RapidAI/RapidOCR")
#             det_model_path = os.path.join(download_path, "onnx", "PP-OCRv5", "det", "ch_PP-OCRv5_server_det.onnx")
#             rec_model_path = os.path.join(download_path, "onnx", "PP-OCRv5", "rec", "ch_PP-OCRv5_rec_server_infer.onnx")
#             cls_model_path = os.path.join(download_path, "onnx", "PP-OCRv4", "cls", "ch_ppocr_mobile_v2.0_cls_infer.onnx")
            
#             ocr_options = RapidOcrOptions(
#                 det_model_path=det_model_path, rec_model_path=rec_model_path, cls_model_path=cls_model_path
#             )

#             pipeline_options = PdfPipelineOptions()
#             pipeline_options.do_ocr = True
#             pipeline_options.do_table_structure = True
#             pipeline_options.table_structure_options.do_cell_matching = True
#             pipeline_options.ocr_options = ocr_options
            
#             converter = DocumentConverter(
#                 format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)}
#             )
        
#         elif pipeline_mode in ['vlm_local', 'vlm_remote']:
#             model_name = ['allenai/olmOCR-2-7B-1025', 'deepseek-ai/DeepSeek-OCR', 'Qwen/Qwen3-Coder-30B-A3B-Instruct', 'google/gemma-3-27b-it', 'Qwen/Qwen3-VL-8B-Instruct', 'Qwen/Qwen2.5-VL-32B-Instruct', 'qwen/qwen2.5-vl-32b-instruct','x-ai/grok-4-fast', 'google/gemini-2.0-flash-001'] #qwen/qwen2.5-vl-32b-instruct
#             vlm_prompt = (
# """You are an expert Document Analyst AI. Your mission is to meticulously analyze an image of a document and convert its **entire content** into a structured Markdown format. You must process the document sequentially from top to bottom, preserving the original order of all elements.

# **Your primary directive is to perform a literal, verbatim extraction. You must not summarize, interpret, rephrase, or omit any content.**

# Your response must be a single, complete Markdown document representing the source file.

# ### Your Guiding Principles

# 1.  **Absolute Sequential Order:** Process the document from its absolute top to its absolute bottom, transcribing elements in the exact order they appear. If content is arranged visually (e.g., side-by-side columns), you must transcribe the content in a logical reading order (e.g., top-to-bottom of the left column, then top-to-bottom of the right column). The order of elements in your output **must** perfectly match the original.

# 2.  **Verbatim Text Transcription:** All text elements must be transcribed **verbatim (exactly as written)** and converted into the appropriate Markdown format. This includes:

#       * Headings (e.g., `# Title`, `## Subtitle`)
#       * Paragraphs (plain text)
#       * Lists (bulleted `*` or numbered `1.`)
#       * Code blocks (using \`\`\` fences)
#       * Bold (`**text**`) and Italic (`*text*`) styling.
#       * Any and all text labels, captions, or annotations, in the exact place they appear.

# 3.  **Literal Deconstruction of Non-Textual Elements (CRITICAL):** This is **not a summary**. This is a **textual conversion** of visual information. For any visual element that cannot be represented as text (such as images, photographs, charts, **circuit diagrams, schematics,** signatures, stamps, or logos), you must:

#       * Provide a **hyper-detailed, component-by-component, literal transcription** of the element's content.
#       * **Zero Summarization:** Your transcription must be exhaustive and deconstruct the element for someone who cannot see it. You must transcribe all labels, all data points, and all connections.
#       * **For Diagrams/Schematics:** This is the most critical task. You must *trace* and *transcribe* **every single connection** between all labeled components. This is a textual representation of the visual data, not a summary of its purpose.
#           * List each component by its label (e.g., "Arduino Pin A0/14/PC0", "Resistor R1", "Dip Switch 1").
#           * Explicitly trace what each pin connects to, including any intermediate components. (e.g., "Arduino Pin 0/PD0 connects to one side of the first switch in 'Dip Switch 1' AND to one end of a '10K-ohm' resistor. The other end of this resistor connects to '+5V'. The other side of the first switch connects to 'GND'.").
#       * **For Charts/Graphs:** State the chart type (bar, line, pie). Transcribe all axes labels, the title, and the legend. Then, list the data points or relationships shown, one by one.
#       * Enclose this entire literal deconstruction within the specific tags defined in Principle 4.

# 4.  **Tagging Format:** You must use the following tags to enclose your non-textual deconstructions. **The text *inside* the tags is the literal transcription of the visual, not a summary.**

#       * **Tables:** Enclose the entire Markdown table within `<table>` and `</table>`.
#         ```markdown
#         <table>
#         | Header 1 | Header 2 |
#         |---|---|
#         | Data 1 | Data 2 |
#         </table>
#         ```
#       * **Charts:**
#         ```markdown
#         <chart>A vertical bar chart titled 'Sales'. The X-axis is 'Month' (Jan, Feb, Mar). The Y-axis is 'Revenue'. The bar for Feb is taller than Jan.</chart>
#         ```
#       * **Images/Photos:**
#         ```markdown
#         <image>A photograph of a white coffee mug with a blue logo, sitting on a wooden desk next to a laptop.</image>
#         ```
#       * **Diagrams/Schematics (Your Most Important Tag):**
#         ```markdown
#         <diagram>
#         A hyper-detailed, connection-by-connection transcription of the schematic.
#         **Component 1 (e.g., Arduino):**
#         * Pin 1 connects to...
#         * Pin 2 connects to Resistor R1...
#         **Component 2 (e.g., 7-Segment Display):**
#         * Pin 'A' connects to the other end of Resistor R1...
#         * Pin 'Common' connects to GND...
#         (This transcription must trace all connections literally from the image.)
#         </diagram>
#         ```
#       * **Logos:**
#         ```markdown
#         <logo>A circular company logo with the text 'Innovate Corp' and a gear icon in the center.</logo>
#         ```
#       * **Signatures:**
#         ```markdown
#         <signature>A handwritten, illegible signature in blue ink.</signature>
#         ```
#       * **Stamps:**
#         ```markdown
#         <stamp>A red circular stamp with the text 'APPROVED' in the center.</stamp>
#         ```

# -----

# --- Example of Correct Approach ---

# ### Gold-Standard Example (Your Target Quality)

# The example you provided is the **perfect** model of this non-summary approach. It correctly provides a literal, connection-by-connection transcription of the visual diagram *first*, and *then* provides a verbatim transcription of all the text that follows it, in the correct order.

# **Your required output must be exactly in this format and at this level of detail:**

# ```markdown
# <diagram>
# A detailed circuit schematic showing connections between an Arduino-style board, a common cathode 7-Segment display, and an 8-position Dip Switch block labeled "ดิปสวิตช์1".

# **Connections from Arduino Board:**

# * **To 7-Segment Display (via 300-ohm resistors):**
#     * Pin A0/14/PC0 connects to a 300-ohm resistor, which connects to the 'A' segment pin.
#     * Pin A1/15/PC1 connects to a 300-ohm resistor, which connects to the 'B' segment pin.
#     * Pin A2/16/PC2 connects to a 300-ohm resistor, which connects to the 'C' segment pin.
#     * Pin A3/17/PC3 connects to a 300-ohm resistor, which connects to the 'D' segment pin.
#     * Pin A4/18/PC4 connects to a 300-ohm resistor, which connects to the 'E' segment pin.
#     * Pin A5/19/PC5 connects to a 300-ohm resistor, which connects to the 'F' segment pin.
#     * Pin 8/PB0 connects to a 300-ohm resistor, which connects to the 'G' segment pin.
# * **To Dip Switch 1 (Pull-Up Configuration):**
#     * Pin 0/PD0 connects to a node. This node is connected to one side of the first switch (leftmost) and also to one end of a 10K-ohm resistor.
#     * Pin 1/PD1 connects to a node. This node is connected to one side of the second switch and also to one end of a 10K-ohm resistor.
#     * Pin 2/PD2 connects to a node. This node is connected to one side of the third switch and also to one end of a 10K-ohm resistor.
#     * Pin 3/PD3 connects to a node. This node is connected to one side of the fourth switch and also to one end of a 10K-ohm resistor.
#     * Pin 4/PD4 connects to a node. This node is connected to one side of the fifth switch and also to one end of a 10K-ohm resistor.
#     * Pin 5/PD5 connects to a node. This node is connected to one side of the sixth switch and also to one end of a 10K-ohm resistor.
#     * Pin 6/PD6 connects to a node. This node is connected to one side of the seventh switch and also to one end of a 10K-ohm resistor.
#     * Pin 7/PD7 connects to a node. This node is connected to one side of the eighth switch (rightmost) and also to one end of a 10K-ohm resistor.

# **Component Connections:**

# * **7-Segment Display:**
#     * Segments A, B, C, D, E, F, and G are connected as described above.
#     * The common cathode pin (at the bottom) is connected to GND.
# * **Dip Switch 1 ("ดิปสวิตช์1"):**
#     * This is an 8-position switch block.
#     * One side of each of the 8 switches is connected to its respective Arduino pin (PD0-PD7) as described above.
#     * The other side of all 8 switches is connected to a common GND line.
# * **Pull-Up Resistors (10K-ohm x 8):**
#     * There are 8 10K-ohm resistors.
#     * One end of each resistor is connected to an Arduino pin node (PD0-PD7).
#     * The other end of all 8 resistors is connected to a common +5V line.

# **Power Pins:**
# * The Arduino board shows a "USB JACK".
# * A power pin block shows 3.3V, 5V, GND, GND, VIN.
# * Unconnected pins on the Arduino board include: SCL, SDA, AREF, GND, 13/PB5, 12/PB4, 11/PB3, 10/PB2, 9/PB1.
# </diagram>

# รูปที่ 1

# Lab 1
# Dip Switch and 7-Segment

# อุปกรณ์
# 1. Arduino Board
# 2. Digital Experiment Board
# 3. 7-Segment Board

# Checkpoint# 1: อ่านค่าสถานะจาก Dip-switch เพื่อแสดงค่าออกที่ 7-Segment
# 1.1 ต่อวงจร ตามรูปที่ 1
# 1.2 ใช้โปรแกรม Arduino IDE เพื่อป้อน Code ด้านล่าง และแก้ไข เพิ่มเติมให้เหมาะสม เพื่อให้ Code สามารถอ่านค่าสถานะตรรกะจาก Dip-Switch ทั้ง 8 แล้วทำการนับจำนวน Switch ที่มี ตรรกะ “HIGH” แล้วแสดงค่าจำนวนนั้นออกสู่ 7-Segment ได้อย่างถูกต้อง
# ```
# --- End of example ---
# """)
#             vlm_mode_arg = 'local' if pipeline_mode == 'vlm_local' else 'remote'
#             pipeline_options = generate_vlm_pipeline_options(mode=vlm_mode_arg, prompt=vlm_prompt, model_name=model_name[6]) # DeepInfra 0 5 OpenRouter 6
#             converter = DocumentConverter(format_options={InputFormat.PDF: PdfFormatOption(pipeline_cls=VlmPipeline, pipeline_options=pipeline_options)})
#         else:
#             return f"Error: Invalid pipeline_mode '{pipeline_mode}'. Choose 'ocr', 'vlm_local', or 'vlm_remote'."
        
#         result = converter.convert(doc_stream)
#         doc = result.document
#     except Exception as e:
#         return f"Error processing file with Docling: {e}"

#     # Get the main markdown content, which includes the '' placeholders
#     full_content = doc.export_to_markdown()

#     # --- Step 3: Get descriptions for the images we extracted earlier ---
#     image_descriptions = []
#     if extracted_images_bytes:
#         print(f"🖼️ Found {len(extracted_images_bytes)} images. Waiting for descriptions...")
#         with concurrent.futures.ThreadPoolExecutor() as executor:
#             # Create a list of futures to maintain order
#             futures = []
#             for img_bytes in extracted_images_bytes:
#                 futures.append(executor.submit(image_to_describe_from_base64, img_bytes))
#                 time.sleep(2) # Slight delay to avoid overwhelming the API
#             for future in concurrent.futures.as_completed(futures):
#                 image_descriptions.append(future.result())
#         print("✅ All image descriptions received.")

#     # --- Step 4 (MODIFIED): Replace placeholders with image descriptions ---
#     if image_descriptions:
#         print("🔄 Replacing image placeholders with descriptions in order...")
#         for description in image_descriptions:
#             # Replace the *first* occurrence of the placeholder tag.
#             # The '1' count is crucial to ensure we replace them one by one sequentially.
#             if '' in full_content:
#                 full_content = full_content.replace(
#                     '',
#                     f"\n\n**--- IMAGE DESCRIPTION ---**\n\n{description}\n\n**--- END IMAGE ---**\n\n",
#                     1
#                 )
#             else:
#                 full_content += f"\n\n**--- IMAGE DESCRIPTION ---**\n\n{description}\n\n**--- END IMAGE ---**\n\n"

#     if option == 'summarize':
#         return summarize_text_with_llm(full_content)
#     else: # 'describe'
#         print("📝 Returning final markdown with inline image descriptions.")
#         return full_content

def open_utf8_file(file) -> str:
    """
    Opens a UTF-8 encoded text file and returns its content as a string.
    """
    try:
        with io.TextIOWrapper(file, encoding='utf-8') as f:
            content = f.read()
        return content
    except Exception as e:
        print(f"❌ Error reading UTF-8 file: {e}")
        return ""


# ==============================================================================
#  LEGACY: EXTRACTOR FUNCTIONS (Unchanged, they call extract_and_process_content)
# ==============================================================================

def read_stream(file_stream: BytesIO, filename: str):
    """
    Reads text from a BytesIO stream based on the filename extension.
    """
    ext = os.path.splitext(filename)[1].lower()
    content = ""

    # Ensure stream is at the beginning
    file_stream.seek(0)

    try:
        # --- Word Processing ---
        if ext == '.docx':
            # python-docx accepts file-like objects directly
            doc = Document(file_stream)
            content = "\n".join([p.text for p in doc.paragraphs])
        
        elif ext == '.odt':
            # odfpy accepts file-like objects directly
            doc = load(file_stream)
            elements = doc.getElementsByType(text.P)
            content = "\n".join([str(e) for e in elements])

        elif ext == '.rtf':
            # RTF is text; decode bytes to string first
            text_content = file_stream.read().decode('utf-8', errors='ignore')
            content = rtf_to_text(text_content)

        # --- Presentations ---
        elif ext == '.pptx':
            # python-pptx accepts file-like objects directly
            prs = Presentation(file_stream)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            content = "\n".join(text_runs)

        # --- Spreadsheets ---
        elif ext in ['.xlsx', '.xlsm']:
            # openpyxl accepts file-like objects directly
            wb = openpyxl.load_workbook(file_stream, data_only=True)
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    content += " ".join([str(cell) for cell in row if cell]) + "\n"

        elif ext == '.xls':
            # xlrd needs the raw bytes passed to 'file_contents'
            wb = xlrd.open_workbook(file_contents=file_stream.getvalue())
            for sheet in wb.sheets():
                for row_idx in range(sheet.nrows):
                    row_vals = sheet.row_values(row_idx)
                    content += " ".join([str(v) for v in row_vals if v]) + "\n"

        # --- Unsupported Legacy Formats ---
        elif ext in ['.doc', '.ppt']:
            return f"[ERROR] .doc and .ppt (Binary) require heavy tools (LibreOffice/Antiword) and cannot be read by lightweight Python scripts."

        else:
            return f"[ERROR] Unsupported extension: {ext}"

    except Exception as e:
        return f"Error reading {filename}: {str(e)}"

    return content

def extract_pdf_text(file_storage, option: str = 'describe', mode: str = 'vlm_remote') -> str:
    """Extracts content from a PDF. Can describe (default) or summarize."""
    return extract_and_process_content(file_storage, option, pipeline_mode=mode, with_images=False)

def extract_image_text(file_storage, option: str = 'describe') -> str:
    """Extracts content from a standalone image. Can describe (default) or summarize."""
    text = image_to_describe_from_base64(file_storage.read())
    print(f"Extracted text from image: {text}")
    return text
    # return process_pages_with_vlm(file_storage.read(),"describe this images")

def extract_txt_file(file, option: str = 'describe', mode: str = 'vlm_remote') -> str:
    """Extracts content from a .txt file. Can describe (default) or summarize."""
    return open_utf8_file(file)

def extract_docx_text(file, option: str = 'describe', mode: str = 'vlm_remote') -> str:
    """Extracts content from a .docx file. Can describe (default) or summarize."""
    return extract_and_process_content(file, option, pipeline_mode=mode, with_images=True)

def extract_pptx_text(file, option: str = 'describe', mode: str = 'vlm_remote') -> str:
    """Extracts content from a .pptx file. Can describe (default) or summarize."""
    return extract_and_process_content(file, option, pipeline_mode=mode, with_images=True)

def extract_excel_text(file, option: str = 'describe', mode: str = 'vlm_remote') -> str:
    """Extracts content from an .xlsx file. Can describe (default) or summarize."""
    return extract_and_process_content(file, option, pipeline_mode=mode, with_images=True)

def extract_xls_text(file, option: str = 'describe', mode: str = 'vlm_remote') -> str:
    """Extracts content from an .xls file. Can describe (default) or summarize."""
    return extract_and_process_content(file, option, pipeline_mode=mode, with_images=True)

def summarize_text_with_llm(text: str) -> str:
    """
    Summarizes a long text using an LLM via OpenRouter.
    """
    print("📚 Summarizing content...")
    system_prompt = "You are an expert summarizer. Your task is to provide a concise and clear summary of the given text, capturing the key points and main ideas."
    prompt = f"Please summarize the following content:\n\n---\n\n{text}\n\n---\n\nSummary:"
    
    # Using a fast and cost-effective model for summarization tasks
    if not LOCAL:
        if IFXGPT:
            IFXGPTInference(prompt=prompt, system_prompt=system_prompt, model_name="gpt-5-mini")
            pass
        else:
            summary = OpenRouterInference(prompt=prompt, system_prompt=system_prompt, model_name="x-ai/grok-4-fast")
    else:
        summary = ollama_generate_text(prompt=prompt, system_prompt=system_prompt, model="gemma3:4b")[0]
    return summary

def image_to_describe_from_base64(image_bytes: bytes) -> str:
    """
    Extract object descriptions from an image in base64 format.
    """
    # save image bytes to a PIL Image
    # image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    # You might want to manage temporary files more carefully in production
    # temp_image_path = f"./temp_images/temp_image_for_description_{uuid.uuid4()}.jpg"
    # os.makedirs(os.path.dirname(temp_image_path), exist_ok=True)
    # image.save(temp_image_path)
    
    # System prompt for OpenRouter VLM
    system_prompt = ("You're an image expert."
                     "If the image contains text, extract and summarize it...")
    # Prompt for OpenRouter VLM
    prompt = ("Please describe the image in detail in a text format that allows you to understand its details.")
    # Call the object detection API with the image bytes
    if not LOCAL:
        if IFXGPT:
            response = IFXGPTInference(prompt=prompt, system_prompt=system_prompt, image_bytes_list=[image_bytes], model_name="gpt-5-mini")
        else:
            response = OpenRouterInference(prompt=prompt, system_prompt=system_prompt, image_bytes_list=[image_bytes], model_name="qwen/qwen2.5-vl-32b-instruct") #qwen/qwen2.5-vl-32b-instruct // qwen/qwen3-vl-8b-instruct
    else:
        response = ollama_describe_image(image_bytes=image_bytes, model="qwen3-vl:2b-instruct", prompt=prompt, system_prompt=system_prompt)
    return response

# ==============================================================================
#  UPDATED: VLM & EMBEDDING INFERENCE
# ==============================================================================

def OpenRouterInference(prompt: str, system_prompt: str = "", image_bytes_list: List[bytes] = None, model_name: str = "google/gemma-3-12b-it") -> str:
    """
    Perform inference using OpenRouter API with optional MULTIPLE image input for VLM.
    This version resizes images and dynamically detects the MIME type.

    Args:
        prompt: The user prompt.
        system_prompt: The system prompt to guide the model.
        image_bytes_list: Optional LIST of image bytes for VLM processing.
        model_name: The name of the OpenRouter model to use.

    Returns:
        The model's response as a string, or an error message.
    """
    api_key = os.getenv("OPENROUTER_API_KEY")
    if not api_key:
        return "Error: OPENROUTER_API_KEY environment variable not set."

    messages = []
    if system_prompt:
        messages.append({"role": "system", "content": system_prompt})

    # Start user content with the text prompt
    user_content = [{"type": "text", "text": prompt}]

    if image_bytes_list:
        print(f"Processing {len(image_bytes_list)} images for VLM...")
        for image_bytes in image_bytes_list:
            try:
                # Open the image from bytes
                img = Image.open(io.BytesIO(image_bytes))
                
                # --- START: RESIZING SECTION ---
                max_size = (1024, 1024)  # Max width and height of 1024 pixels
                img.thumbnail(max_size, Image.Resampling.LANCZOS)
                
                output_buffer = io.BytesIO()
                image_format = img.format or 'JPEG' 
                img.save(output_buffer, format=image_format)
                
                resized_image_bytes = output_buffer.getvalue()
                # --- END: RESIZING SECTION ---

                base64_image = base64.b64encode(resized_image_bytes).decode('utf-8')
                image_url = f"data:image/{image_format.lower()};base64,{base64_image}"
                
                # Append each image to the user_content list
                user_content.append({
                    "type": "image_url",
                    "image_url": {"url": image_url}
                })
            except Exception as e:
                print(f"Error processing, resizing, or encoding an image: {e}")
                # Optionally add a text placeholder for the failed image
                user_content.append({"type": "text", "text": "[Image processing failed]"})
    
    messages.append({"role": "user", "content": user_content})

    # --- The rest of the function remains the same ---
    try:
        response = requests.post(
            url="https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            },
            json={
                "model": model_name,
                "messages": messages
            }
        )
        response.raise_for_status()
        data = response.json()
        return data['choices'][0]['message']['content']

    except requests.exceptions.RequestException as e:
        return f"Error calling OpenRouter API: {e}"
    except (KeyError, IndexError) as e:
        return f"Error parsing OpenRouter response: {e}. Full response: {response.text}"
    except Exception as e:
        return f"An unexpected error occurred: {e}"



# def extract_and_process_content(file_storage, option: str = 'describe', pipeline_mode: str = 'ocr', with_images: bool = True) -> str:
#     """
#     UPDATED: This function now processes files by converting them to images 
#     and sending them to a VLM for analysis.
    
#     - For PDFs: Converts each page to an image and sends all page images.
#     - For Images: Sends the single image.
#     - Other types (DOCX, TXT, etc.) fall back to the legacy OCR text extraction.
#     """
    
#     # --- Step 1: Get file bytes ---
#     if isinstance(file_storage, bytes):
#         file_bytes = file_storage
#         file_storage = io.BytesIO(file_bytes)
#         file_storage.filename = "temp_file.pdf" # Mock filename
#     else:
#         file_storage.seek(0)
#         file_bytes = file_storage.read()
#         file_storage.seek(0) # Reset pointer
        
#     if not file_storage.filename:
#         # Try to guess type if no filename
#         file_storage.filename = "temp_file.bin"

#     file_ext = os.path.splitext(file_storage.filename)[1].lower()
#     image_bytes_list = []
    
#     # --- Step 2: Convert file to list of images ---
    
#     if file_ext == '.pdf':
#         print(f"Processing PDF '{file_storage.filename}': Converting pages to images...")
#         try:
#             pdf_document = fitz.open(stream=file_bytes, filetype="pdf")
#             page_count = len(pdf_document)
#             print(f"PDF has {page_count} pages.")
            
#             for page_num_0_idx in range(page_count):
#                 page_image_bytes = convert_pdf_page_to_image(file_bytes, page_num_0_idx, dpi=200)
#                 if page_image_bytes:
#                     image_bytes_list.append(page_image_bytes)
#                 else:
#                     print(f"⚠️ Warning: Could not convert page {page_num_0_idx + 1}")
            
#             pdf_document.close()
#             print(f"✅ Successfully converted {len(image_bytes_list)} pages to images.")
            
#         except Exception as e:
#             return f"Error opening or converting PDF: {e}"
            

#     elif file_ext in ['.png', '.jpg', '.jpeg', '.webp', '.bmp', '.gif']:
#         print(f"Processing single image: '{file_storage.filename}'")
#         image_bytes_list.append(file_bytes)
        
#     else:
#         # Fallback for unsupported types: Try the old OCR logic
#         # This preserves functionality for DOCX, TXT, etc.
#         print(f"Unsupported file type '{file_ext}' for VLM-image flow. Falling back to legacy OCR...")
#         # --- ORIGINAL OCR/TEXT EXTRACTION LOGIC (Simplified) ---
#         try:
#             doc_stream = DocumentStream(name=file_storage.filename, stream=io.BytesIO(file_bytes))
#             # print("Downloading RapidOCR models")
#             # download_path = snapshot_download(repo_id="RapidAI/RapidOCR")
#             # det_model_path = os.path.join(download_path, "onnx", "PP-OCRv5", "det", "ch_PP-OCRv5_server_det.onnx")
#             # rec_model_path = os.path.join(download_path, "onnx", "PP-OCRv5", "rec", "ch_PP-OCRv5_rec_server_infer.onnx")
#             # cls_model_path = os.path.join(download_path, "onnx", "PP-OCRv4", "cls", "ch_ppocr_mobile_v2.0_cls_infer.onnx")
            
#             # ocr_options = RapidOcrOptions(
#             #     det_model_path=det_model_path, rec_model_path=rec_model_path, cls_model_path=cls_model_path
#             # )
#             # pipeline_options = PdfPipelineOptions()
#             # pipeline_options.do_ocr = True
#             # pipeline_options.do_table_structure = True
#             # pipeline_options.table_structure_options.do_cell_matching = True
#             # pipeline_options.ocr_options = ocr_options

#             # pipeline_options_word = WordPipelineOptions()
#             # pipeline_options_powerpoint = PowerPointPipelineOptions()
#             # pipeline_options_excel = ExcelPipelineOptions()
            
#             # Use a generic converter for fallback
#             converter = DocumentConverter(
#                 # format_options={
#                 #     InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options),
#                 #     InputFormat.DOCX: PdfFormatOption(pipeline_options=pipeline_options_word), 
#                 #     InputFormat.PPTX: PdfFormatOption(pipeline_options=pipeline_options_powerpoint),
#                 #     InputFormat.XLSX: PdfFormatOption(pipeline_options=pipeline_options_excel),
#                 #     # InputFormat.XLS: PdfFormatOption(pipeline_options=pipeline_options),
#                 #     # InputFormat.TXT: PdfFormatOption(pipeline_options=pipeline_options),
#                 # }
#             )
            
#             result = converter.convert(doc_stream)
#             doc = result.document
#             full_content = doc.export_to_markdown()
            
#             if option == 'summarize':
#                 return summarize_text_with_llm(full_content)
#             else:
#                 return full_content
                
#         except Exception as e:
#             return f"Error during legacy OCR fallback: {e}"
#         # --- END OF FALLBACK ---

#     if not image_bytes_list:
#         return "Error: No images were successfully extracted or converted from the file."

#     # --- Step 3: Define VLM Prompt ---
#     # (This is the detailed instruction prompt from your original code)
#     vlm_system_prompt = (
# """You are an expert Document Analyst AI. Your mission is to meticulously analyze an image of a document and convert its **entire content** into a structured Markdown format. You must process the document sequentially from top to bottom, preserving the original order of all elements.

# **Your primary directive is to perform a literal, verbatim extraction. You must not summarize, interpret, rephrase, or omit any content.**

# Your response must be a single, complete Markdown document representing the source file.

# ### Your Guiding Principles

# 1.  **Absolute Sequential Order:** Process the document from its absolute top to its absolute bottom, transcribing elements in the exact order they appear. If content is arranged visually (e.g., side-by-side columns), you must transcribe the content in a logical reading order (e.g., top-to-bottom of the left column, then top-to-bottom of the right column). The order of elements in your output **must** perfectly match the original.

# 2.  **Verbatim Text Transcription:** All text elements must be transcribed **verbatim (exactly as written)** and converted into the appropriate Markdown format. This includes:

#       * Headings (e.g., `# Title`, `## Subtitle`)
#       * Paragraphs (plain text)
#       * Lists (bulleted `*` or numbered `1.`)
#       * Code blocks (using \`\`\` fences)
#       * Bold (`**text**`) and Italic (`*text*`) styling.
#       * Any and all text labels, captions, or annotations, in the exact place they appear.

# 3.  **Literal Deconstruction of Non-Textual Elements (CRITICAL):** This is **not a summary**. This is a **textual conversion** of visual information. For any visual element that cannot be represented as text (such as images, photographs, charts, **circuit diagrams, schematics,** signatures, stamps, or logos), you must:

#       * Provide a **hyper-detailed, component-by-component, literal transcription** of the element's content.
#       * **Zero Summarization:** Your transcription must be exhaustive and deconstruct the element for someone who cannot see it. You must transcribe all labels, all data points, and all connections.
#       * **For Diagrams/Schematics:** This is the most critical task. You must *trace* and *transcribe* **every single connection** between all labeled components. This is a textual representation of the visual data, not a summary of its purpose.
#           * List each component by its label (e.g., "Arduino Pin A0/14/PC0", "Resistor R1", "Dip Switch 1").
#           * Explicitly trace what each pin connects to, including any intermediate components. (e.g., "Arduino Pin 0/PD0 connects to one side of the first switch in 'Dip Switch 1' AND to one end of a '10K-ohm' resistor. The other end of this resistor connects to '+5V'. The other side of the first switch connects to 'GND'.").
#       * **For Charts/Graphs:** State the chart type (bar, line, pie). Transcribe all axes labels, the title, and the legend. Then, list the data points or relationships shown, one by one.
#       * Enclose this entire literal deconstruction within the specific tags defined in Principle 4.

# 4.  **Tagging Format:** You must use the following tags to enclose your non-textual deconstructions. **The text *inside* the tags is the literal transcription of the visual, not a summary.**

#       * **Tables:** Enclose the entire Markdown table within `<table>` and `</table>`.
#         ```markdown
#         <table>
#         | Header 1 | Header 2 |
#         |---|---|
#         | Data 1 | Data 2 |
#         </table>
#         ```
#       * **Charts:**
#         ```markdown
#         <chart>A vertical bar chart titled 'Sales'. The X-axis is 'Month' (Jan, Feb, Mar). The Y-axis is 'Revenue'. The bar for Feb is taller than Jan.</chart>
#         ```
#       * **Images/Photos:**
#         ```markdown
#         <image>A photograph of a white coffee mug with a blue logo, sitting on a wooden desk next to a laptop.</image>
#         ```
#       * **Diagrams/Schematics (Your Most Important Tag):**
#         ```markdown
#         <diagram>
#         A hyper-detailed, connection-by-connection transcription of the schematic.
#         **Component 1 (e.g., Arduino):**
#         * Pin 1 connects to...
#         * Pin 2 connects to Resistor R1...
#         **Component 2 (e.g., 7-Segment Display):**
#         * Pin 'A' connects to the other end of Resistor R1...
#         * Pin 'Common' connects to GND...
#         (This transcription must trace all connections literally from the image.)
#         </diagram>
#         ```
#       * **Logos:**
#         ```markdown
#         <logo>A circular company logo with the text 'Innovate Corp' and a gear icon in the center.</logo>
#         ```
#       * **Signatures:**
#         ```markdown
#         <signature>A handwritten, illegible signature in blue ink.</signature>
#         ```
#       * **Stamps:**
#         ```markdown
#         <stamp>A red circular stamp with the text 'APPROVED' in the center.</stamp>
#         ```

# -----

# ### *** Note ***
# You must extract all content from the document.

# ### Gold-Standard Example (Your Target Quality)

# The example you provided is the **perfect** model of this non-summary approach. It correctly provides a literal, connection-by-connection transcription of the visual diagram *first*, and *then* provides a verbatim transcription of all the text that follows it, in the correct order.

# **Your required output must be exactly in this format and at this level of detail:**

# ```markdown
# <diagram>
# A detailed circuit schematic showing connections between an Arduino-style board, a common cathode 7-Segment display, and an 8-position Dip Switch block labeled "ดิปสวิตช์1".

# **Connections from Arduino Board:**

# * **To 7-Segment Display (via 300-ohm resistors):**
#     * Pin A0/14/PC0 connects to a 300-ohm resistor, which connects to the 'A' segment pin.
#     * Pin A1/15/PC1 connects to a 300-ohm resistor, which connects to the 'B' segment pin.
#     * Pin A2/16/PC2 connects to a 300-ohm resistor, which connects to the 'C' segment pin.
#     * Pin A3/17/PC3 connects to a 300-ohm resistor, which connects to the 'D' segment pin.
#     * Pin A4/18/PC4 connects to a 300-ohm resistor, which connects to the 'E' segment pin.
#     * Pin A5/19/PC5 connects to a 300-ohm resistor, which connects to the 'F' segment pin.
#     * Pin 8/PB0 connects to a 300-ohm resistor, which connects to the 'G' segment pin.
# * **To Dip Switch 1 (Pull-Up Configuration):**
#     * Pin 0/PD0 connects to a node. This node is connected to one side of the first switch (leftmost) and also to one end of a 10K-ohm resistor.
#     * Pin 1/PD1 connects to a node. This node is connected to one side of the second switch and also to one end of a 10K-ohm resistor.
#     * Pin 2/PD2 connects to a node. This node is connected to one side of the third switch and also to one end of a 10K-ohm resistor.
#     * Pin 3/PD3 connects to a node. This node is connected to one side of the fourth switch and also to one end of a 10K-ohm resistor.
#     * Pin 4/PD4 connects to a node. This node is connected to one side of the fifth switch and also to one end of a 10K-ohm resistor.
#     * Pin 5/PD5 connects to a node. This node is connected to one side of the sixth switch and also to one end of a 10K-ohm resistor.
#     * Pin 6/PD6 connects to a node. This node is connected to one side of the seventh switch and also to one end of a 10K-ohm resistor.
#     * Pin 7/PD7 connects to a node. This node is connected to one side of the eighth switch (rightmost) and also to one end of a 10K-ohm resistor.

# **Component Connections:**

# * **7-Segment Display:**
#     * Segments A, B, C, D, E, F, and G are connected as described above.
#     * The common cathode pin (at the bottom) is connected to GND.
# * **Dip Switch 1 ("ดิปสวิตช์1"):**
#     * This is an 8-position switch block.
#     * One side of each of the 8 switches is connected to its respective Arduino pin (PD0-PD7) as described above.
#     * The other side of all 8 switches is connected to a common GND line.
# * **Pull-Up Resistors (10K-ohm x 8):**
#     * There are 8 10K-ohm resistors.
#     * One end of each resistor is connected to an Arduino pin node (PD0-PD7).
#     * The other end of all 8 resistors is connected to a common +5V line.

# **Power Pins:**
# * The Arduino board shows a "USB JACK".
# * A power pin block shows 3.3V, 5V, GND, GND, VIN.
# * Unconnected pins on the Arduino board include: SCL, SDA, AREF, GND, 13/PB5, 12/PB4, 11/PB3, 10/PB2, 9/PB1.
# </diagram>

# รูปที่ 1

# Lab 1
# Dip Switch and 7-Segment

# อุปกรณ์
# 1. Arduino Board
# 2. Digital Experiment Board
# 3. 7-Segment Board

# Checkpoint# 1: อ่านค่าสถานะจาก Dip-switch เพื่อแสดงค่าออกที่ 7-Segment
# 1.1 ต่อวงจร ตามรูปที่ 1
# 1.2 ใช้โปรแกรม Arduino IDE เพื่อป้อน Code ด้านล่าง และแก้ไข เพิ่มเติมให้เหมาะสม เพื่อให้ Code สามารถอ่านค่าสถานะตรรกะจาก Dip-Switch ทั้ง 8 แล้วทำการนับจำนวน Switch ที่มี ตรรกะ “HIGH” แล้วแสดงค่าจำนวนนั้นออกสู่ 7-Segment ได้อย่างถูกต้อง
# ```""")
    
#     vlm_system_prompt = (
# """You are an expert Document Analyst AI converting images to structured Markdown.

# **CORE DIRECTIVE: Extract content sequentially and verbatim. NO summarization, interpretation, or omission.**

# ### Operational Rules
# 1.  **Sequential Order:** Transcribe elements top-to-bottom, left-to-right.
# 2.  **Text Transcription:** Extract text exactly as written, preserving Markdown formatting (Headers, Lists, Code blocks, **Bold**, *Italic*).
# 3.  **Visual Deconstruction (CRITICAL):**
#     *   Convert visuals into literal text descriptions inside specific tags.
#     *   **For Schematics/Diagrams:** You must provide a **hyper-detailed, pin-by-pin connection trace**. Explicitly state every wire connection (e.g., "Pin A connects to Resistor R1, which connects to GND"). Describe the structure, not the function.

# ### Required Tags
# *   `<diagram>`: Detailed schematic connection tracing.
# *   `<table>`: Markdown tables.
# *   `<chart>`: Type, axes, legend, and data points.
# *   `<image>` / `<logo>` / `<signature>` / `<stamp>`: Literal visual description.

# ### Relevance Filter
# If a specific user question is provided and this page contains no relevant information to answer it, do not extract content of that section.
# """)
#     # Add context about the images
#     page_references = [f"- Document Page {i+1}" for i in range(len(image_bytes_list))]
#     reference_text = "\n".join(page_references)

#     # This prompt is sent as the 'user' message
#     final_user_prompt = f"""

#     I am providing you with {len(image_bytes_list)} images from the file '{file_storage.filename}'. These images represent the pages of the document in sequential order: {reference_text}

#     Please analyze all these pages as a single, continuous document and generate the full Markdown extraction as requested in the system prompt. Begin processing from Page 1 and continue sequentially to the end. """

#     # --- Step 4: Call OpenRouterInference ---
#     print(f"Sending {len(image_bytes_list)} images to DeepInfraInference VLM...")

#     vlm_response = ""
#     batch_size = 10
#     for i in range(len(image_bytes_list)//batch_size + 2): # Process in batches of 15
#         print(f"Processing images {(i)*batch_size + 1} to {min((i+1)*batch_size, len(image_bytes_list))}...")
#         if (i)*batch_size >= len(image_bytes_list):
#             break

#         if not LOCAL:
#             if IFXGPT:
#                 vlm_response += IFXGPTInference(
#                     prompt=final_user_prompt,
#                     system_prompt=vlm_system_prompt, # The user's detailed instructions go here
#                     image_bytes_list=image_bytes_list[(i)*batch_size:(i+1)*batch_size], # Send in batches of 15 images
#                     model_name= "gpt-5-mini" #"Qwen/Qwen3-VL-8B-Instruct" #"deepseek-ai/DeepSeek-OCR" #"Qwen/Qwen2.5-VL-32B-Instruct" # Using a strong VLM Qwen/Qwen3-VL-8B-Instruct Qwen/Qwen3-VL-30B-A3B-Instruct Qwen/Qwen2.5-VL-32B-Instruct
#                 ) + "\n\n"
#             else:
#                 vlm_response += DeepInfraInference(
#                     prompt=final_user_prompt,
#                     system_prompt=vlm_system_prompt, # The user's detailed instructions go here
#                     image_bytes_list=image_bytes_list[(i)*batch_size:(i+1)*batch_size], # Send in batches of 15 images
#                     model_name= "meta-llama/Llama-4-Scout-17B-16E-Instruct" #"Qwen/Qwen3-VL-8B-Instruct" #"deepseek-ai/DeepSeek-OCR" #"Qwen/Qwen2.5-VL-32B-Instruct" # Using a strong VLM Qwen/Qwen3-VL-8B-Instruct Qwen/Qwen3-VL-30B-A3B-Instruct Qwen/Qwen2.5-VL-32B-Instruct
#                 ) + "\n\n"
#             print(vlm_response)
#         else :
#              # System prompt for OpenRouter VLM
#             vlm_system_prompt = ("You're an image expert."
#                              "If the image contains text, extract and summarize it...")
#             # Prompt for OpenRouter VLM
#             final_user_prompt = ("Please describe the image in detail in a text format that allows you to understand its details.")
#             vlm_response_L = ollama_describe_image(
#                 image_bytes=image_bytes_list[(i)*batch_size:(i+1)*batch_size],
#                 model="qwen3-vl:2b-instruct",
#                 prompt=final_user_prompt,
#                 system_prompt=vlm_system_prompt
#                 )
#             vlm_response += "\n\n".join(vlm_response_L) + "\n\n"

#     print(vlm_response)
#     print("✅ VLM processing complete.")

#     # --- Step 5: Handle "summarize" option ---
#     if option == 'summarize':
#         print("Summarizing VLM output...")
#         return summarize_text_with_llm(vlm_response)
#     else: # 'describe'
#         return vlm_response


def extract_and_process_content(file_storage, option: str = 'describe', pipeline_mode: str = 'ocr', with_images: bool = True) -> str:
    """
    UPDATED (Memory Optimized): Processes files by converting them to images 
    and sending them to a VLM for analysis.
    
    - For PDFs: Processes in BATCHES to save memory (does not load all pages at once).
    - For Images: Sends the single image.
    - Other types: Fallback to legacy OCR.
    """
    import gc  # Garbage collector for explicit memory management

    # --- Step 1: Get file bytes ---
    if isinstance(file_storage, bytes):
        file_bytes = file_storage
        file_storage = io.BytesIO(file_bytes)
        file_storage.filename = "temp_file.pdf"
    else:
        file_storage.seek(0)
        file_bytes = file_storage.read()
        file_storage.seek(0)

    if not file_storage.filename:
        file_storage.filename = "temp_file.bin"

    file_ext = os.path.splitext(file_storage.filename)[1].lower()
    
    # Define batch size (tune this based on your VRAM/RAM limits)
    BATCH_SIZE = 10 
    vlm_response = ""

    # --- Step 2: Define VLM Prompt (Static Part) ---
    vlm_system_prompt = (
"""You are an expert Document Analyst AI converting images to structured Markdown.

**CORE DIRECTIVE: Extract content sequentially and verbatim. NO summarization, interpretation, or omission.**

### Operational Rules
1.  **Sequential Order:** Transcribe elements top-to-bottom, left-to-right.
2.  **Text Transcription:** Extract text exactly as written, preserving Markdown formatting.
3.  **Visual Deconstruction:** Convert visuals into literal text descriptions inside specific tags (<diagram>, <table>, <chart>, <image>).
4.  **Schematics:** Provide hyper-detailed, pin-by-pin connection traces.

### Relevance Filter
If a specific user question is provided and this page contains no relevant information, do not extract content of that section.
""")

    # --- Step 3: Process based on file type ---

    if file_ext == '.pdf':
        print(f"Processing PDF '{file_storage.filename}' in batches of {BATCH_SIZE}...")
        try:
            pdf_document = fitz.open(stream=file_bytes, filetype="pdf")
            page_count = len(pdf_document)
            print(f"PDF has {page_count} pages.")
            
            # Loop through the PDF in chunks
            for start_idx in range(0, page_count, BATCH_SIZE):
                end_idx = min(start_idx + BATCH_SIZE, page_count)
                print(f"Processing Batch: Pages {start_idx + 1} to {end_idx}...")
                
                batch_images = []
                page_references = []

                # Convert only the current batch of pages to images
                for page_num in range(start_idx, end_idx):
                    # convert_pdf_page_to_image must return bytes
                    page_image_bytes = convert_pdf_page_to_image(file_bytes, page_num, dpi=200)
                    if page_image_bytes:
                        batch_images.append(page_image_bytes)
                        page_references.append(f"- Document Page {page_num + 1}")
                    else:
                        print(f"⚠️ Warning: Could not convert page {page_num + 1}")

                if not batch_images:
                    continue

                # Prepare dynamic prompt for this batch
                reference_text = "\n".join(page_references)
                batch_user_prompt = f"""
                I am providing you with {len(batch_images)} images representing a section of the document: 
                {reference_text}
                
                Please analyze these pages and extract the content as requested. 
                """

                # Call VLM for this batch
                if not LOCAL:
                    if IFXGPT:
                        batch_res = IFXGPTInference(
                            prompt=batch_user_prompt,
                            system_prompt=vlm_system_prompt,
                            image_bytes_list=batch_images,
                            model_name="gpt-5-mini" 
                        )
                    else:
                        batch_res = DeepInfraInference(
                            prompt=batch_user_prompt,
                            system_prompt=vlm_system_prompt,
                            image_bytes_list=batch_images,
                            model_name="meta-llama/Llama-4-Scout-17B-16E-Instruct"
                        )
                    vlm_response += batch_res + "\n\n"
                else:
                    # Local Ollama fallback
                    desc_list = ollama_describe_image(
                        image_bytes=batch_images,
                        model="qwen3-vl:2b-instruct",
                        prompt="Please describe the image in detail.",
                        system_prompt="You are an image expert."
                    )
                    vlm_response += "\n\n".join(desc_list) + "\n\n"

                # CRITICAL: Clear memory for this batch
                del batch_images
                gc.collect() 

            pdf_document.close()
            print("✅ PDF processing complete.")

        except Exception as e:
            return f"Error opening or converting PDF: {e}"

    elif file_ext in ['.png', '.jpg', '.jpeg', '.webp', '.bmp', '.gif']:
        # Single image path (same as before but simplified)
        print(f"Processing single image: '{file_storage.filename}'")
        image_bytes_list = [file_bytes]
        
        user_prompt = "Analyze this image and extract content."
        
        if not LOCAL:
             # Call your inference function (simplified for brevity)
             if IFXGPT:
                vlm_response = IFXGPTInference(prompt=user_prompt, system_prompt=vlm_system_prompt, image_bytes_list=image_bytes_list, model_name="gpt-5-mini")
             else:
                vlm_response = DeepInfraInference(prompt=user_prompt, system_prompt=vlm_system_prompt, image_bytes_list=image_bytes_list, model_name="meta-llama/Llama-4-Scout-17B-16E-Instruct")
        else:
             res_list = ollama_describe_image(image_bytes=image_bytes_list, model="qwen3-vl:2b-instruct", prompt="Describe", system_prompt="Expert")
             vlm_response = res_list[0]

    else:
        # Fallback for DOCX/TXT/etc using Legacy OCR
        print(f"Unsupported file type '{file_ext}' for VLM-image flow. Falling back to legacy OCR...")
        try:
            doc_stream = DocumentStream(name=file_storage.filename, stream=io.BytesIO(file_bytes))
            converter = DocumentConverter() # Your legacy converter setup
            result = converter.convert(doc_stream)
            vlm_response = result.document.export_to_markdown()
        except Exception as e:
            return f"Error during legacy OCR fallback: {e}"

    # --- Step 4: Final Output ---
    if option == 'summarize':
        print("Summarizing VLM output...")
        return summarize_text_with_llm(vlm_response)
    else:
        return vlm_response
    
    
# ==============================================================================
#  OLLAMA FUNCTION (NEW)
# ==============================================================================

def OllamaInference(prompt: str, system_prompt: str = "", image_bytes_list: List[bytes] = None, model_name: str = "llava") -> str:
    """
    Perform inference using Ollama's local API.
    This version supports optional MULTIPLE image input for VLM (models like llava).
    
    Args:
        prompt: The user prompt.
        system_prompt: The system prompt to guide the model.
        image_bytes_list: Optional LIST of image bytes for VLM processing.
        model_name: The name of the Ollama model to use (e.g., 'llava', 'llama2-vision').
    
    Returns:
        The model's response as a string, or an error message.
    """
    ollama_host = os.getenv("OLLAMA_HOST", "http://localhost:11434")
    api_url = f"{ollama_host}/api/generate"
    
    # Prepare the prompt with images if provided
    full_prompt = prompt
    
    if image_bytes_list:
        print(f"Processing {len(image_bytes_list)} images for Ollama VLM...")
        # For Ollama, we can embed base64 images in the prompt
        for idx, image_bytes in enumerate(image_bytes_list):
            try:
                base64_image = base64.b64encode(image_bytes).decode('utf-8')
                full_prompt += f"\n[Image {idx + 1}]: (base64 image data attached)"
            except Exception as e:
                print(f"Error encoding image {idx}: {e}")
                full_prompt += f"\n[Image {idx + 1} processing failed]"
    
    # Add system prompt context
    if system_prompt:
        full_prompt = system_prompt + "\n\n" + full_prompt
    
    try:
        # Ollama API request
        response = requests.post(
            api_url,
            json={
                "model": model_name,
                "prompt": full_prompt,
                "stream": False,
                "temperature": 0.0,
            },
            timeout=300  # 5 minute timeout for complex queries
        )
        response.raise_for_status()
        data = response.json()
        
        # Ollama returns response in 'response' field
        return data.get('response', 'No response from Ollama')
    
    except requests.exceptions.ConnectionError:
        return f"Error: Could not connect to Ollama at {ollama_host}. Make sure Ollama is running."
    except requests.exceptions.RequestException as e:
        return f"Error calling Ollama API: {e}"
    except (KeyError, IndexError) as e:
        return f"Error parsing Ollama response: {e}. Full response: {response.text if response else 'No response'}"
    except Exception as e:
        return f"An unexpected error occurred with Ollama: {e}"


# ==============================================================================
#  DEEPINFRA FUNCTION (Updated)
# ==============================================================================

parameter_option = {
                'normal' : {
                    'temperature': 0.1,
                    'top_p': 0.95,
                    "frequency_penalty": 0.0,
                    "presence_penalty": 0.0,
                },

                'deepseek-ai/DeepSeek-OCR' : {
                    'temperature': 0.1,
                    'top_p': 0.9,
                    'min_p': 0.0,
                    "frequency_penalty": 0.0,
                    "presence_penalty": 0.0,
                },
            }

def DeepInfraInference(prompt: str = "", system_prompt: str = "", image_bytes_list: List[bytes] = None, model_name: str = "deepseek-ai/DeepSeek-OCR") -> str:
    """
    Perform inference using DeepInfra's OpenAI-compatible API.
    This version supports optional MULTIPLE image input for VLM, resizing,
    and returns output in the same format as OpenRouter.

    Args:
        prompt: The user prompt.
        system_prompt: The system prompt to guide the model.
        image_bytes_list: Optional LIST of image bytes for VLM processing.
        model_name: The name of the DeepInfra model to use.

    Returns:
        The model's response as a string, or an error message.
    """
    api_key = os.getenv("DEEPINFRA_API_KEY")
    if not api_key:
        return "Error: DEEPINFRA_API_KEY environment variable not set."

    messages = []
    if system_prompt:
        messages.append({"role": "system", "content": system_prompt})

    # Start user content with the text prompt
    user_content = [{"type": "text", "text": prompt}]

    if image_bytes_list:
        print(f"Processing {len(image_bytes_list)} images for VLM...")
        for image_bytes in image_bytes_list:
            try:
                # Open the image from bytes
                img = Image.open(io.BytesIO(image_bytes))
                
                # --- START: RESIZING SECTION ---
                # max_size = (1028, 1028)  # Max width and height of 1028 pixels
                # img.thumbnail(max_size, Image.Resampling.LANCZOS)
                
                output_buffer = io.BytesIO()
                image_format = img.format or 'JPEG' 
                img.save(output_buffer, format=image_format)
                
                resized_image_bytes = output_buffer.getvalue()
                # --- END: RESIZING SECTION ---

                base64_image = base64.b64encode(resized_image_bytes).decode('utf-8')
                image_url = f"data:image/{image_format.lower()};base64,{base64_image}"
                
                # Append each image to the user_content list
                user_content.append({
                    "type": "image_url",
                    "image_url": {"url": image_url}
                })
            except Exception as e:
                print(f"Error processing, resizing, or encoding an image: {e}")
                # Optionally add a text placeholder for the failed image
                user_content.append({"type": "text", "text": "[Image processing failed]"})
    
    messages.append({"role": "user", "content": user_content})

    # --- API Call Section (Modified for DeepInfra) ---
    try:
        if model_name in parameter_option.keys():
            parameter = parameter_option[model_name]
        else:
            parameter = parameter_option['normal']
        parameter['model'] = model_name
        parameter['messages'] = messages
        response = requests.post(
            # Use DeepInfra's OpenAI-compatible endpoint
            url="https://api.deepinfra.com/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            },
            json=parameter
        )
        response.raise_for_status()
        data = response.json()
        # Parse the response in the same way as OpenRouter
        return data['choices'][0]['message']['content']

    except requests.exceptions.RequestException as e:
        return f"Error calling DeepInfra API: {e}"
    except (KeyError, IndexError) as e:
        return f"Error parsing DeepInfra response: {e}. Full response: {response.text}"
    except Exception as e:
        return f"An unexpected error occurred: {e}"

def DeepInfraEmbedding(inputs: list[str], model_name: str = "Qwen/Qwen3-Embedding-0.6B", dimensions: int = 1024) -> list[list[float]]:
    """
    Generate embeddings using the DeepInfra API.
    (Original function, unchanged)
    """
    api_key = os.getenv("DEEPINFRA_API_KEY")
    if not api_key:
        print("Error: DEEPINFRA_API_KEY environment variable not set.")
        return []

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }

    payload = {"inputs": inputs, "dimensions": dimensions}
    
    api_url = f"https://api.deepinfra.com/v1/inference/{model_name}"

    try:
        response = requests.post(url=api_url, headers=headers, json=payload)
        response.raise_for_status()
        data = response.json()
        return data.get('embeddings', []) # Safely get the embeddings

    except requests.exceptions.RequestException as e:
        print(f"Error calling DeepInfra Embedding API: {e}")
        return []
    except (KeyError, IndexError) as e:
        print(f"Error parsing DeepInfra Embedding response: {e}. Full response: {response.text}")
        return []
    except Exception as e:
        print(f"An unexpected error occurred during embedding: {e}")
        return []

# --- NEW: DeepInfra CLIP Embedding Function ---
# def get_image_embedding_jinna_api(text: str = None, image_bytes: bytes = None, model_name: str = "sentence-transformers/clip-ViT-B-32") -> Optional[List[float]]:
#     """
#     Gets a CLIP embedding from DeepInfra for either text or an image.
    
#     Args:
#         text: The text string to embed.
#         image_bytes: The raw image bytes to embed.
#         model_name: The CLIP model on DeepInfra.

#     Returns:
#         A 512-dimension embedding vector, or None on failure.
#     """
#     api_key = os.getenv("DEEPINFRA_API_KEY")
#     if not api_key:
#         print("Error: DEEPINFRA_API_KEY environment variable not set.")
#         return None
#     if not text and not image_bytes:
#         print("Error: Must provide either text or image_bytes to get_image_embedding_jinna_api.")
#         return None

#     headers = {"Authorization": f"Bearer {api_key}"}
#     api_url = f"https://api.deepinfra.com/v1/inference/{model_name}"
    
#     payload = {}
#     if text:
#         payload = {"inputs": [text]}
#     elif image_bytes:
#         # Convert image bytes to base64 data URL
#         base64_img = base64.b64encode(image_bytes).decode('utf-8')
#         # Guess mime type, default to jpeg if unknown
#         mime_type = "image/jpeg"
#         try:
#             img = Image.open(io.BytesIO(image_bytes))
#             if img.format:
#                 mime_type = Image.MIME.get(img.format)
#         except Exception:
#             pass # Use default
        
#         data_url = f"data:{mime_type};base64,{base64_img}"
#         payload = {"image": data_url}

#     try:
#         print("Requesting CLIP embedding from DeepInfra...")
#         response = requests.post(url=api_url, headers=headers, json=payload)
#         response.raise_for_status()
#         data = response.json()
        
#         embeddings = data.get('embeddings', [])
#         if embeddings and len(embeddings) > 0:
#             print(f"✅ Generated CLIP embedding using DeepInfra (Type: {'Text' if text else 'Image'}).")
#             return embeddings[0]
#         else:
#             print(f"❌ DeepInfra CLIP API response did not contain embeddings: {data}")
#             return None

#     except requests.exceptions.RequestException as e:
#         print(f"Error calling DeepInfra CLIP API: {e}")
#         return None
#     except Exception as e:
#         print(f"An unexpected error occurred during CLIP embedding: {e}")
#         return None


# --- UPDATED: Jina v4 Multimodal Embedding Function ---
def get_image_embedding_jinna_api(
    text: str = None, 
    search_text: str = None,
    image_bytes_list: List[bytes] = None, 
    model_name: str = "jina-embeddings-v4"
) -> Union[Optional[List[float]], Optional[List[List[float]]]]:
    """
    Gets a multimodal embedding from Jina AI v4.

    - If 'text' is provided, embeds a single text string and returns one embedding (List[float]).
    - If 'image_bytes_list' is provided, embeds a batch of images and returns a list of embeddings (List[List[float]]).
    
    Args:
        text: The text string to embed.
        image_bytes_list: A list of raw image bytes to embed.
        model_name: The Jina model to use.

    Returns:
        A single embedding vector (List[float]) if 'text' was used.
        A list of embedding vectors (List[List[float]]) if 'image_bytes_list' was used.
        None on failure.
    """
    api_key = os.getenv("JINA_API_KEY")
    if not api_key:
        print("Error: JINA_API_KEY environment variable not set.")
        return None
    
    # Ensure only one input type is provided
    if (text or search_text) and image_bytes_list:
        print("Error: Provide either 'text' OR 'image_bytes_list', not both.")
        return None
    if not (text or search_text) and not image_bytes_list:
        print("Error: Must provide either 'text' or 'image_bytes_list' to get_image_embedding_jinna_api.")
        return None

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}" # The API key from Jina dashboard
    }
    api_url = "https://api.jina.ai/v1/embeddings"
    
    input_data = []
    task_type = ""
    # create_search_prompt = f"Help me to create prompt for use to search page by text from this task or question add more detail to explain the document page. \n\n task or question : {text} \n\n Give me only new prompt no additional text."
#     create_search_prompt = f"""
# Act as a document search that simurch engine. 
# Based on the user's query below, generate a detailed paragraph describing the content, specific keywords, and technical terminology likely to appear on a document page that answers this query in english. 
# Do not answer the question directly; only describe the page content.

# User Query: {text}

# Output only the descriptive paragraph. No introductory text.
# """


#     create_search_prompt = f"""
# Act as a document search engine (PDF document search by vector similarity). 
# Write a single, concise sentence that simulates a direct excerpt from a document page answering the query below. 
# Include likely keywords and factual phrasing.

# User Query: {text}
# Type of Document: Datasheet or Manual (Table, Graph, Diagram or Text)
# Prompt Language: English

# Output only the simulated excerpt.
# """ #*****************
    







    create_search_prompt = f"""
Act as a high-precision OCR and document retrieval system.
Generate a single, dense excerpt that mimics a direct "hit" from a technical manual or datasheet answering the query: "{text}".
Requirements:
    1. Use high-density technical language (e.g., "nominal operating range," "tolerance ±5%," "hex-head bolt," "thermal conductivity").
    2. Include specific identifiers like model numbers (e.g., Series-X1), section headers (e.g., Section 4.2.1), or units (e.g., mV, PSI, Nm).
    3. Structure the output as a standalone technical fact or a row from a specification table.
    4. Do not use introductory filler; start immediately with the data.
    """
    
    # create_search_prompt = f"""
# Act as a document search engine. Simulate a raw text extraction from a table or diagram within a PDF datasheet to answer: "{text}".
# Style Guidelines:
#     # Format the response as a condensed data string (e.g., [Parameter] | [Value] | [Conditions]).
#     # Use "telegraphic" writing style (no unnecessary articles like 'the' or 'a').
#     # Ensure the text contains a specific numerical value or a direct instruction found in a "Notes" column.
#     # Example format: "Table 3: Max Load Capacity; Model A: 500kg; Model B: 750kg; Note: Ensure surface is level."
# """

    # create_search_prompt = f"""
# Act as a semantic search engine retrieving a snippet from a field service manual.
# Generate one concise, authoritative sentence that answers the query "{text}" as it would appear in a printed instruction set.
# Instructions:
#     # Use imperative verbs (e.g., "Calibrate," "Verify," "Torque," "Connect").
#     # Include a page or figure reference (e.g., "See Fig. 12-B").
#     # Use specific technical terminology relevant to the query to ensure high vector similarity.
#     # Output only the excerpt.
# """









#     create_search_prompt = f"""
# Act as a document retrieval expert. 
# Analyze the User Query and generate a concise search string to find the relevant document page. 
# Include specific technical terminology, synonyms, and key concepts that would appear in the text of the target page.
# Keep the output short, dense, and strictly relevant.

# User Query: {text}

# Output only the search string.
# """
    
    if text or search_text:
        print("Requesting Jina v4 embedding (Type: Text)...")
        if search_text == None:
            if not LOCAL:
                if IFXGPT:
                    search_text = IFXGPTInference(
                        prompt=create_search_prompt,
                        # system_prompt=system_prompt,
                        # image_bytes_list=image_bytes_list,
                        model_name="gpt-5-mini" #'x-ai/grok-4-fast'#"Qwen/Qwen2.5-VL-32B-Instruct"#"Qwen/Qwen3-235B-A22B-Instruct-2507" # Use a strong VLM
                    )
                else:
                    search_text = DeepInfraInference(
                        prompt=create_search_prompt,
                        # system_prompt=system_prompt,
                        # image_bytes_list=image_bytes_list,
                        model_name="Qwen/Qwen3-235B-A22B-Instruct-2507" #'x-ai/grok-4-fast'#"Qwen/Qwen2.5-VL-32B-Instruct"#"Qwen/Qwen3-235B-A22B-Instruct-2507" # Use a strong VLM
                    )

            else :
                search_text = ollama_generate_text(
                    prompt=create_search_prompt,
                    model="gemma3:4b"
                )
        else:
            search_text = search_text

        print(f"Search prompt: {search_text}")
        input_data.append({"text": search_text})
        task_type = "retrieval.query"
    
    elif image_bytes_list:
        print(f"Requesting Jina v4 embedding (Type: {len(image_bytes_list)} Images)...")
        for img_bytes in image_bytes_list:
            # Jina API accepts a raw base64 string for the image
            base64_img = base64.b64encode(img_bytes).decode('utf-8')
            input_data.append({"image": base64_img})
            task_type = "retrieval.passage"

    payload = {
        "model": model_name,
        "task": task_type, # Use "text-matching" as shown in the API example for similarity
        "input": input_data,
        "dimensions": 2048
    }

    try:
        # Batch requests when sending many images to Jina to avoid large single requests / SSL errors.
        aggregated_data = []
        try:
            # If images were provided and input is large, send in batches
            if image_bytes_list and len(input_data) > 100:
                batch_size = 100
                offset = 0
                resp = None
                import random
                max_retries = 8
                for start in range(0, len(input_data), batch_size):
                    chunk = input_data[start:start + batch_size]
                    print(len(chunk))
                    chunk_payload = {
                        "model": model_name,
                        "task": task_type,
                        "input": chunk,
                        # "dimensions": 256,
                    }
                    retries = 0
                    while retries < max_retries:
                        try:
                            resp = requests.post(url=api_url, headers=headers, json=chunk_payload, timeout=120)
                            if resp.status_code == 429:
                                wait_time = (2 ** retries) + random.uniform(0, 1)
                                print(f"429 Too Many Requests. Retrying batch {start//batch_size+1} after {wait_time:.2f}s (attempt {retries+1}/{max_retries})...")
                                time.sleep(wait_time)
                                retries += 1
                                continue
                            resp.raise_for_status()
                            chunk_json = resp.json()
                            chunk_data = chunk_json.get("data", [])

                            # adjust indices so they reflect global position across batches
                            for item in chunk_data:
                                item["index"] = item.get("index", 0) + offset
                            print(len(chunk_data))
                            aggregated_data.extend(chunk_data)
                            print(len(aggregated_data))
                            offset += len(chunk)
                            print(f"batch :{((start+1)//batch_size) + 1}/{len(input_data)//batch_size} is finish")
                            time.sleep(1)
                            break
                        except requests.exceptions.RequestException as e:
                            wait_time = (2 ** retries) + random.uniform(0, 1)
                            print(f"Request error: {e}. Retrying batch {start//batch_size+1} after {wait_time:.2f}s (attempt {retries+1}/{max_retries})...")
                            time.sleep(wait_time)
                            retries += 1
                    else:
                        print(f"❌ Failed to get embeddings for batch {start//batch_size+1} after {max_retries} retries.")
                # emulate single combined response structure expected later
                data = {"data": aggregated_data}
                response = resp
            else:
                # single request (text or small image batch)
                max_retries = 8
                retries = 0
                import random
                while retries < max_retries:
                    response = requests.post(url=api_url, headers=headers, json=payload, timeout=120)
                    if response.status_code == 429:
                        wait_time = (2 ** retries) + random.uniform(0, 1)
                        print(f"429 Too Many Requests. Retrying after {wait_time:.2f}s (attempt {retries+1}/{max_retries})...")
                        time.sleep(wait_time)
                        retries += 1
                        continue
                    try:
                        response.raise_for_status()
                        data = response.json()
                        break
                    except requests.exceptions.RequestException as e:
                        wait_time = (2 ** retries) + random.uniform(0, 1)
                        print(f"Request error: {e}. Retrying after {wait_time:.2f}s (attempt {retries+1}/{max_retries})...")
                        time.sleep(wait_time)
                        retries += 1
                else:
                    print(f"❌ Failed to get embeddings after {max_retries} retries.")
                    return None
                response.raise_for_status()
                data = response.json()
        # except requests.exceptions.RequestException as e:
            # print(f"Error calling Jina Embedding API: {e}")
            # return None
        except Exception as e:
            print(f"Unexpected error during Jina embedding requests: {e}")
            return None
        
        # Parse the Jina response format
        embedding_data = data.get('data', [])
        if not embedding_data:
            print(f"❌ Jina API response did not contain data: {data}")
            return None

        # Sort by index to ensure order is preserved
        sorted_data = sorted(embedding_data, key=lambda x: x['index'])
        embeddings_list = [item['embedding'] for item in sorted_data]

        if not embeddings_list:
             print(f"❌ Jina API response did not contain valid embeddings: {data}")
             return None

        # Return a single list if text was the input (for search_similar_pages)
        if text or search_text:
            print(f"✅ Generated Jina v4 embedding for text.")
            return embeddings_list[0] 
        
        # Return a list of lists if images were the input (for /process)
        elif image_bytes_list:
            print(f"✅ Generated {len(embeddings_list)} Jina v4 embeddings for images.")
            return embeddings_list

    except requests.exceptions.RequestException as e:
        print(f"Error calling Jina Embedding API: {e}")
        return None
    except (KeyError, IndexError) as e:
        # This catches errors if the response structure is not as expected
        print(f"Error parsing Jina Embedding response: {e}. Full response: {response.text}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred during Jina embedding: {e}")
        return None

# Global cache for the model to prevent reloading it on every function call (High Latency)
# _JINA_MODEL_INSTANCE = None

def get_image_embedding_jinna_api_local(
    text: str = None, 
    search_text: str = None,
    image_bytes_list: List[bytes] = None, 
    model_name: str = "jinaai/jina-embeddings-v4"
) -> Union[Optional[List[float]], Optional[List[List[float]]]]:
    """
    Gets a multimodal embedding locally using SentenceTransformer (Jina v4).

    - If 'text' is provided: Generates a hypothetical document description using an LLM, 
      then embeds that description. Returns one embedding (List[float]).
    - If 'image_bytes_list' is provided: Embeds a batch of images. 
      Returns a list of embeddings (List[List[float]]).
    
    Args:
        text: The text string to embed.
        image_bytes_list: A list of raw image bytes to embed.
        model_name: The HuggingFace model ID to use.

    Returns:
        A single embedding vector (List[float]) if 'text' was used.
        A list of embedding vectors (List[List[float]]) if 'image_bytes_list' was used.
        None on failure.
    """
    global _JINA_MODEL_INSTANCE

    # 1. Input Validation
    if (text or search_text) and image_bytes_list:
        print("Error: Provide either 'text' OR 'image_bytes_list', not both.")
        return None
    if not text and not search_text and not image_bytes_list:
        print("Error: Must provide either 'text' or 'image_bytes_list'.")
        return None

    try:
        # 2. Load Model (Singleton Pattern)
        # if _JINA_MODEL_INSTANCE is None:
        #     print(f"Loading local model: {model_name}...")
        #     _JINA_MODEL_INSTANCE = SentenceTransformer(model_name, trust_remote_code=True)
        #     print("Model loaded successfully.")
        
        # model = _JINA_MODEL_INSTANCE

        # 3. Handle Text Input (Retrieval Query)
        if text or search_text:
            print("Generating Jina v4 embedding (Type: Text)...")
            
            # --- START: Query Expansion / HyDE Logic (Preserved) ---
            create_search_prompt = f"""
Act as a document search engine. 
Based on the user's query below, generate a detailed paragraph describing the content, specific keywords, and technical terminology likely to appear on a document page that answers this query. 
Do not answer the question directly; only describe the page content.

User Query: {text}

Output only the descriptive paragraph. No introductory text.
"""
            if search_text == None:
                # Ensure LOCAL and helper functions are defined in your outer scope
                if not LOCAL:
                    if IFXGPT:
                        search_text = IFXGPTInference(
                            prompt=create_search_prompt,
                            model_name="gpt-5-mini" 
                        )
                    else:
                        search_text = DeepInfraInference(
                            prompt=create_search_prompt,
                            model_name="Qwen/Qwen3-235B-A22B-Instruct-2507" 
                        )
                else:
                    search_text = ollama_generate_text(
                        prompt=create_search_prompt,
                        model="gemma3:4b"
                    )
            else:
                search_text = search_text
            print(f"Search prompt (HyDE): {search_text}")
            # --- END: Query Expansion Logic ---

            # Encode text
            # Task 'retrieval.query' optimizes the embedding for finding matching documents
            # model.to(device)
            with torch.no_grad():
                embedding = model.encode([search_text], task="retrieval",convert_to_numpy=True)
                # model.to("cpu")
                clear_gpu()
            
            # Convert numpy array to list
            print(f"✅ Generated Jina v4 embedding for text.")
            return embedding[0].tolist()
        
        # 4. Handle Image Input (Retrieval Passage)
        elif image_bytes_list:
            print(f"Generating Jina v4 embedding (Type: {len(image_bytes_list)} Images)...")
            
            # Convert raw bytes to PIL Images
            pil_images = []
            for img_bytes in image_bytes_list:
                pil_images.append(Image.open(io.BytesIO(img_bytes)))
            
            # Encode images (Batch processing is handled automatically by SentenceTransformer)
            # Task 'retrieval.passage' optimizes the embedding for being indexed
            # Note: Ensure the specific Jina model supports image inputs (like Jina-CLIP or specific V4 variants)
            # model.to(device)
            with torch.no_grad():
                embeddings = model.encode(pil_images,batch_size=1, convert_to_numpy=True) # task="retrieval.passage" is implied for non-query inputs usually, or add if model supports
                # model.to("cpu")
                clear_gpu()
            
            print(f"✅ Generated {len(embeddings)} Jina v4 embeddings for images.")
            return embeddings.tolist()

    except Exception as e:
        print(f"An unexpected error occurred during local embedding generation: {e}")
        return None


def get_image_embedding_local_api_colpali_engine(
    text: str = None, 
    image_bytes_list: List[bytes] = None, 
    model_name: str = "vidore/colSmol-250M"
) -> Union[Optional[List[float]], Optional[List[List[float]]]]:
    """
    Gets embeddings using the local ColPali engine.

    - If 'text' is provided, embeds a single text string and returns one embedding (List[float]).
    - If 'image_bytes_list' is provided, embeds a batch of images and returns a list of embeddings (List[List[float]]).
    
    Args:
        text: The text string to embed.
        image_bytes_list: A list of raw image bytes to embed.
        model_name: The ColPali model name (default: "vidore/colSmol-500M").

    Returns:
        A single embedding vector (List[float]) if 'text' was used.
        A list of embedding vectors (List[List[float]]) if 'image_bytes_list' was used.
        None on failure.
    """
    # Ensure only one input type is provided
    if text and image_bytes_list:
        print("Error: Provide either 'text' OR 'image_bytes_list', not both.")
        return None
    if not text and not image_bytes_list:
        print("Error: Must provide either 'text' or 'image_bytes_list' to get_image_embedding_local_api_colpali_engine.")
        return None

    try:
        # Load the model and processor
        model = ColIdefics3.from_pretrained(
            model_name,
            torch_dtype=torch.float16,
            device_map=device,
            attn_implementation="flash_attention_2"
        ).eval()
        print(f"Loaded ColPali model: {model_name} on device: {model.device}")
        processor = ColIdefics3Processor.from_pretrained(model_name)

        if text:
            print("Requesting ColPali embedding (Type: Text)...")
            # Process the text query
            batch_queries = processor.process_queries([text]).to(model.device)
            with torch.no_grad():
                query_embeddings = model(**batch_queries)
            # Assuming query_embeddings is a tensor, convert to list
            embedding = query_embeddings.cpu().numpy().tolist()
            if embedding:
                print("✅ Generated ColPali embedding for text.")
                return embedding[0]  # Single embedding
            else:
                print("❌ Failed to generate embedding for text.")
                return None
        
        elif image_bytes_list:
            print(f"Requesting ColPali embedding (Type: {len(image_bytes_list)} Images)...")
            # Convert image bytes to PIL Images
            images = []
            for img_bytes in image_bytes_list:
                try:
                    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                    images.append(img)
                except Exception as e:
                    print(f"Error processing image: {e}")
                    images.append(Image.new("RGB", (32, 32), color="white"))  # Placeholder
            
            # Process the images
            # Batch process images in groups of 3 for efficiency and memory management
            batch_embeddings = []
            batch_size = 3
            for i in range(0, len(images), batch_size):
                batch_imgs = images[i:i + batch_size]
                batch_images = processor.process_images(batch_imgs).to(model.device)
                with torch.no_grad():
                    image_embeddings = model(**batch_images)
                    print(image_embeddings.shape)
                    time.sleep(10)
                    # Pad embeddings from 128 to 1024 dimensions (FIXED from 2048 to match other embeddings)
                    if image_embeddings.shape[1] < 1024:
                        pad_width = 1024 - image_embeddings.shape[1]
                        image_embeddings = torch.nn.functional.pad(image_embeddings, (0, pad_width), mode='constant', value=0)
                    elif image_embeddings.shape[1] > 1024:
                        image_embeddings = image_embeddings[:, :1024]
                    batch_embeddings.extend(image_embeddings.cpu().numpy().tolist())
            embeddings = batch_embeddings
            # with torch.no_grad():
            #     image_embeddings = model(**batch_images)
            #     # Pad embeddings from 128 to 1024 dimensions (FIXED from 2048)
            #     if image_embeddings.shape[1] < 1024:
            #         pad_width = 1024 - image_embeddings.shape[1]
            #         image_embeddings = torch.nn.functional.pad(image_embeddings, (0, pad_width), mode='constant', value=0)
            #     elif image_embeddings.shape[1] > 1024:
            #         image_embeddings = image_embeddings[:, :1024]
            # Convert to list of lists
            # embeddings = image_embeddings.cpu().numpy().tolist()
            if embeddings:
                print(f"✅ Generated {len(embeddings)} ColPali embeddings for images.")
                return embeddings
            else:
                print("❌ Failed to generate embeddings for images.")
                return None

    except Exception as e:
        print(f"Error in ColPali embedding: {e}")
        return None


# def get_image_embedding_jinna_api(
#     text: str = None, 
#     image_bytes_list: List[bytes] = None, 
#     model_name: str = "jina-embeddings-v4"
# ) -> Union[Optional[List[float]], Optional[List[List[float]]]]:
#     """
#     Gets a multimodal embedding from Jina AI v4.

#     - If 'text' is provided, embeds a single text string and returns one embedding (List[float]).
#     - If 'image_bytes_list' is provided, embeds a batch of images and returns a list of embeddings (List[List[float]]).
    
#     Args:
#         text: The text string to embed.
#         image_bytes_list: A list of raw image bytes to embed.
#         model_name: The Jina model to use.

#     Returns:
#         A single embedding vector (List[float]) if 'text' was used.
#         A list of embedding vectors (List[List[float]]) if 'image_bytes_list' was used.
#         None on failure.
#     """
#     api_key = os.getenv("JINA_API_KEY")
#     if not api_key:
#         print("Error: JINA_API_KEY environment variable not set.")
#         return None
    
#     # Ensure only one input type is provided
#     if text and image_bytes_list:
#         print("Error: Provide either 'text' OR 'image_bytes_list', not both.")
#         return None
#     if not text and not image_bytes_list:
#         print("Error: Must provide either 'text' or 'image_bytes_list' to get_image_embedding_jinna_api.")
#         return None

#     headers = {
#         "Content-Type": "application/json",
#         "Authorization": f"Bearer {api_key}" # The API key from Jina dashboard
#     }
#     api_url = "https://api.jina.ai/v1/embeddings"
    
#     input_data = []
#     if text:
#         print("Requesting Jina v4 embedding (Type: Text)...")
#         input_data.append(text)
    
#     elif image_bytes_list:
#         print(f"Requesting Jina v4 embedding (Type: {len(image_bytes_list)} Images)...")
#         for img_bytes in image_bytes_list:
#             # Convert image bytes to base64 string before passing to the model
#             base64_img = base64.b64encode(img_bytes).decode('utf-8')
#             input_data.append(base64_img)

#     try:
#         # Batch requests when sending many images to Jina to avoid large single requests / SSL errors.
#         aggregated_data = []
#         try:
#             batch_size = 8
#             embeddings_list = model.encode(sentences=input_data, batch_size=batch_size, show_progress_bar=True, device='cpu', convert_to_numpy=True, task='retrieval', precision="int8")
       
#         except Exception as e:
#             print(f"Unexpected error during Jina embedding requests: {e}")
#             return None

#         if not embeddings_list:
#              print(f"❌ Jina API response did not contain valid embeddings: {None}")
#              return None

#         # Return a single list if text was the input (for search_similar_pages)
#         if text:
#             print(f"✅ Generated Jina v4 embedding for text.")
#             return embeddings_list[0] 
        
#         # Return a list of lists if images were the input (for /process)
#         elif image_bytes_list:
#             print(f"✅ Generated {len(embeddings_list)} Jina v4 embeddings for images.")
#             return embeddings_list

#     except requests.exceptions.RequestException as e:
#         print(f"Error calling Jina Embedding API: {e}")
#         return None
#     except (KeyError, IndexError) as e:
#         # This catches errors if the response structure is not as expected
#         print(f"Error parsing Jina Embedding response: {e}. Full response: {response.text}")
#         return None
#     except Exception as e:
#         print(f"An unexpected error occurred during Jina embedding: {e}")
#         return None

# _JINA_MODEL_INSTANCE = None

# def get_image_embedding_nemoretriever_api_local(
#     text: str = None, 
#     image_bytes_list: List[bytes] = None, 
#     model_name: str = "nvidia/llama-nemoretriever-colembed-1b-v1"
# ) -> Union[Optional[List[float]], Optional[List[List[float]]]]:
#     """
#     Gets embeddings using the NVIDIA Llama Nemoretriever Colembed model.

#     - If 'text' is provided, generates a hypothetical document description using an LLM (HyDE), 
#       then embeds that description as a query. Returns one embedding (List[float]).
#     - If 'image_bytes_list' is provided, embeds a batch of images as passages. 
#       Returns a list of embeddings (List[List[float]]).
    
#     Args:
#         text: The text string to embed.
#         image_bytes_list: A list of raw image bytes to embed.
#         model_name: The HuggingFace model ID (default: NVIDIA model).

#     Returns:
#         A single embedding vector (List[float]) if 'text' was used.
#         A list of embedding vectors (List[List[float]]) if 'image_bytes_list' was used.
#         None on failure.
#     """
#     # global _JINA_MODEL_INSTANCE
#     clear_gpu()

#     # 1. Input Validation
#     if text and image_bytes_list:
#         print("Error: Provide either 'text' OR 'image_bytes_list', not both.")
#         return None
#     if not text and not image_bytes_list:
#         print("Error: Must provide either 'text' or 'image_bytes_list'.")
#         return None

#     try:
#         # 2. Load Model (Singleton Pattern)
#         # if _JINA_MODEL_INSTANCE is None:
#         #     print(f"Loading local model: {model_name}...")
#         #     _JINA_MODEL_INSTANCE = AutoModel.from_pretrained(
#         #         model_name,
#         #         device_map='cuda' if torch.cuda.is_available() else 'cpu',
#         #         trust_remote_code=True,
#         #         torch_dtype=torch.float8_e4m3fn,
#         #         attn_implementation="flash_attention_2",
#         #         revision='1f0fdea7f5b19532a750be109b19072d719b8177'
#         #     ).eval()
#         #     print("Model loaded successfully.")
        
#         # model = _JINA_MODEL_INSTANCE

#         # 3. Handle Text Input (Query Embedding with HyDE)
#         if text:
#             print("Generating query embedding (Type: Text) with HyDE...")
            
#             # --- START: Query Expansion / HyDE Logic ---
#             create_search_prompt = f"""
# Act as a document search engine. 
# Based on the user's query below, generate a detailed paragraph describing the content, specific keywords, and technical terminology likely to appear on a document page that answers this query. 
# Do not answer the question directly; only describe the page content.

# User Query: {text}

# Output only the descriptive paragraph. No introductory text.
# """
#             # Ensure LOCAL and helper functions are defined in your outer scope
#             if not LOCAL:
#                 search_text = DeepInfraInference(
#                     prompt=create_search_prompt,
#                     model_name="Qwen/Qwen3-235B-A22B-Instruct-2507" 
#                 )
#             else:
#                 search_text = ollama_generate_text(
#                     prompt=create_search_prompt,
#                     model="gemma3:1b"
#                 )[0]
            
#             print(f"Search prompt (HyDE): {search_text}")
#             # --- END: Query Expansion Logic ---

#             # Embed the search text as a query
#             queries = [search_text]
#             query_embeddings = model.forward_queries(queries, batch_size=8)
#             # Assuming query_embeddings is a tensor, convert to list
#             # The model outputs multi-vector representations; we take the mean or flatten as needed
#             # For simplicity, flatten or average to a single vector per query
#             embedding = query_embeddings.mean(dim=1).cpu().numpy().tolist()[0]  # Adjust based on model output shape
            
#             print(f" Generated query embedding for text.")
#             return embedding
        
#         # 4. Handle Image Input (Passage Embedding)
#         elif image_bytes_list:
#             print(f"Generating passage embedding (Type: {len(image_bytes_list)} Images)...")
            
#             # Convert raw bytes to PIL Images
#             images = []
#             for img_bytes in image_bytes_list:
#                 try:
#                     img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
#                     original_width, original_height = img.size
#                     new_height = 640  # Example: desired new height
#                     aspect_ratio = original_width / original_height
#                     new_width = int(new_height * aspect_ratio)
#                     new_size_aspect_ratio = (new_width, new_height)
#                     resized_image_aspect_ratio = img.resize(new_size_aspect_ratio)
#                     images.append(resized_image_aspect_ratio)
#                 except Exception as e:
#                     print(f"Error processing image: {e}")
#                     images.append(Image.new("RGB", (224, 224), color="white"))  # Placeholder
            
#             # Embed images as passages
#             # passage_embeddings = model.forward_passages(images, batch_size=1)
#             # Manual inference loop: process images in batches of 3
#             batch_size_manual = 1
#             all_embeddings = []
#             for i in range(0, len(images), batch_size_manual):
#                 batch_images = images[i:i + batch_size_manual]
#                 print(f"Processing batch {i//batch_size_manual + 1} with {len(batch_images)} images...")
                
#                 # Call forward_passages for this batch
#                 batch_embeddings = model(batch_images)
                
#                 # Handle if batch_embeddings is a list (unlikely for small batch, but safe)
#                 if isinstance(batch_embeddings, list):
#                     batch_embeddings = torch.cat(batch_embeddings, dim=0)
                
#                 # batch_embeddings shape: [batch_size, seq_len, embed_dim]
#                 # Average over seq_len to get [batch_size, embed_dim]
#                 batch_embeddings_avg = batch_embeddings.mean(dim=1).cpu().numpy().tolist()
#                 all_embeddings.extend(batch_embeddings_avg)

#             print("finish")
#             print(len(all_embeddings))
#             print(len(all_embeddings[0]))

#             # Check for count mismatch
#             if len(all_embeddings) != len(image_bytes_list):
#                 print(f"- FAILED to get embeddings or count mismatch. Expected {len(image_bytes_list)}, Got {len(all_embeddings)}.")
#                 return None
            
#             print(f" Generated {len(all_embeddings)} passage embeddings for images.")
#             return all_embeddings

#     except Exception as e:
#         print(f"An unexpected error occurred during local embedding generation: {e}")
#         return None

def get_image_embedding_nemoretriever_api_local(
    text: str = None, 
    image_bytes_list: List[bytes] = None, 
    model_name: str = "nvidia/llama-nemoretriever-colembed-1b-v1"
) -> Union[Optional[List[float]], Optional[List[List[float]]]]:
    """
    Gets embeddings using the NVIDIA Llama Nemoretriever Colembed model with manual forward.
    - For text: Use HyDE to generate description, embed as query.
    - For images: Manual forward with positional tensor for passages.
    """
    clear_gpu()

    if text and image_bytes_list:
        print("Error: Provide either 'text' OR 'image_bytes_list', not both.")
        return None
    if not text and not image_bytes_list:
        print("Error: Must provide either 'text' or 'image_bytes_list'.")
        return None

    device = next(model.parameters()).device  # Use model's device

    try:
        if text:
            print("Generating query embedding (Type: Text) with HyDE...")
            # HyDE logic (unchanged)
            create_search_prompt = f"""
Act as a document search engine. 
Based on the user's query below, generate a detailed paragraph describing the content, specific keywords, and technical terminology likely to appear on a document page that answers this query. 
Do not answer the question directly; only describe the page content.

User Query: {text}

Output only the descriptive paragraph. No introductory text.
"""
            if not LOCAL:
                search_text = DeepInfraInference(prompt=create_search_prompt, model_name="Qwen/Qwen3-235B-A22B-Instruct-2507")
            else:
                search_text = ollama_generate_text(prompt=create_search_prompt, model="gemma3:4b")[0]
            print(f"Search prompt (HyDE): {search_text}")

            # For text queries, use language_model forward (manual or model.forward_texts if available)
            # Assuming model has forward_queries; if not, tokenize manually
            from transformers import AutoTokenizer
            tokenizer = AutoTokenizer.from_pretrained(model_name, trust_remote_code=True)
            inputs = tokenizer(search_text, return_tensors="pt", padding=True, truncation=True).to(device)
            with torch.no_grad():
                query_emb = model.language_model(**inputs).last_hidden_state.mean(dim=1)  # Pool to [1, 2048]
            embedding = query_emb.cpu().numpy().tolist()[0]
            print(f"Generated query embedding for text (dim: {len(embedding)}).")
            return embedding

        elif image_bytes_list:
            print(f"Generating passage embedding (Type: {len(image_bytes_list)} Images)...")
            
            # Convert bytes to PIL Images and resize/normalize
            images = []
            for img_bytes in image_bytes_list:
                try:
                    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                    # Resize to model's expected size (512x512 for 1024 patches)
                    img = img.resize((512, 512), Image.Resampling.LANCZOS)
                    images.append(img)
                except Exception as e:
                    print(f"Error processing image: {e}")
                    images.append(Image.new("RGB", (512, 512), color="white"))  # Placeholder

            # Batch process (small batches to avoid OOM)
            batch_size = 1  # Adjust based on GPU memory
            all_embeddings = []
            for i in range(0, len(images), batch_size):
                batch_images = images[i:i + batch_size]
                print(f"Processing batch {i//batch_size + 1} with {len(batch_images)} images...")

                # Step 2a: Preprocess to pixel_values [batch, 3, 512, 512]
                # Provide dummy text input to satisfy processor requirements
                inputs = processor(
                                    images=batch_images,
                                    text=[""] * len(batch_images),
                                    return_tensors="pt",
                                    padding=True,          # Pad text to max length
                                    truncation=True,       # Truncate if too long
                                    # max_length=77,         # Set max text length
                                    do_resize=True,        # Resize images
                                    size={"height": 512, "width": 512},  # Target image size
                                    do_normalize=True      # Normalize images
                                ).to(device)
                pixel_values = inputs['pixel_values']  # [batch, 3, 512, 512]

                # Step 2b: Manual Vision Embeddings with Positional Tensor
                vision_emb = model.vision_model.embeddings
                batch_size_b = pixel_values.shape[0]

                # Patch embeddings: [batch, 1536, 32, 32]
                patch_emb = vision_emb.patch_embedding(pixel_values)

                # Flatten patches: [batch, 1536, 1024] -> [batch, 1024, 1536]
                batch_size_p, embed_dim_p, h_p, w_p = patch_emb.shape
                flattened_patches = patch_emb.flatten(2).transpose(1, 2)  # [batch, 1024, 1536]

                # Create Positional Tensor: Indices for 1024 patches
                # This is the key: torch.arange(0, 1024) for positions 0 to 1023 (row-major order)
                positional_indices = torch.arange(1024, dtype=torch.long, device=device)
                positional_indices = positional_indices.unsqueeze(0).expand(batch_size_b, -1)  # [batch, 1024]

                # Get positional embeddings: [batch, 1024, 1536]
                pos_emb = vision_emb.position_embedding(positional_indices)

                # Add to patches: [batch, 1024, 1536]
                embedded_patches = flattened_patches + pos_emb

                # Step 2c: Forward through vision encoder [batch, 1024, 1536] -> [batch, 1024, 1536]
                encoder_inputs = {'inputs_embeds': embedded_patches}  # Bypass class token if none
                vision_features = model.vision_model.encoder(**encoder_inputs)

                # Post-LN: [batch, 1024, 1536]
                vision_features = model.vision_model.post_layernorm(vision_features.last_hidden_state)

                # Step 2d: Pool via head (MultiheadAttentionPoolingHead) to [batch, 1536]
                pooled_features = model.vision_model.head(vision_features).pooled_output  # Assume this exists; adjust if needed (e.g., mean pool)

                # Step 2e: Project to 2048-dim via mlp1? (mlp1 takes 6144, perhaps concat or multi-vector)
                # For simplicity, mean pool vision_features if head doesn't pool, then project
                # If model expects fusion, use full model forward
                if hasattr(model, 'forward_passages'):
                    # If custom method exists, use it
                    batch_emb = model.forward_passages(batch_images)
                else:
                    # Full model forward (preferred if available)
                    with torch.no_grad():
                        full_outputs = model(pixel_values=pixel_values, output_hidden_states=True)
                        # Extract vision part; assume last_hidden_state is [batch, seq, 2048]
                        batch_emb = full_outputs.last_hidden_state.mean(dim=1)  # Pool to [batch, 2048]

                # Convert to list
                batch_emb_list = batch_emb.cpu().numpy().tolist()
                all_embeddings.extend(batch_emb_list)

            # Check count
            if len(all_embeddings) != len(image_bytes_list):
                print(f"Count mismatch. Expected {len(image_bytes_list)}, got {len(all_embeddings)}.")
                return None

            print(f"Generated {len(all_embeddings)} passage embeddings for images (dim: {len(all_embeddings[0]) if all_embeddings else 0}).")
            return all_embeddings

    except Exception as e:
        print(f"Error during Nemoretriever embedding: {e}")
        import traceback
        traceback.print_exc()
        return None
    

# Global vLLM instance for text generation (HyDE)

# def get_image_embedding_jinna_api_local_vllm(
#     text: str = None, 
#     image_bytes_list: List[bytes] = None, 
#     model_name: str = "jinaai/jina-embeddings-v4"
# ) -> Union[Optional[List[float]], Optional[List[List[float]]]]:
#     """
#     Gets a multimodal embedding locally using SentenceTransformer (Jina v4), but uses vLLM for HyDE text generation.

#     - If 'text' is provided: Generates a hypothetical document description using vLLM (HyDE), 
#       then embeds that description. Returns one embedding (List[float]).
#     - If 'image_bytes_list' is provided: Embeds a batch of images. 
#       Returns a list of embeddings (List[List[float]]).
    
#     Args:
#         text: The text string to embed.
#         image_bytes_list: A list of raw image bytes to embed.
#         model_name: The HuggingFace model ID to use.

#     Returns:
#         A single embedding vector (List[float]) if 'text' was used.
#         A list of embedding vectors (List[List[float]]) if 'image_bytes_list' was used.
#         None on failure.
#     """

#     # 1. Input Validation
#     if text and image_bytes_list:
#         print("Error: Provide either 'text' OR 'image_bytes_list', not both.")
#         return None
#     if not text and not image_bytes_list:
#         print("Error: Must provide either 'text' or 'image_bytes_list'.")
#         return None

#     try:
#         if text:
#             print("Generating Jina v4 embedding (Type: Text) with HyDE using vLLM...")
            
#             # --- START: Query Expansion / HyDE Logic (Updated to use vLLM) ---
#             create_search_prompt = f"""
# Act as a document search engine. 
# Based on the user's query below, generate a detailed paragraph describing the content, specific keywords, and technical terminology likely to appear on a document page that answers this query. 
# Do not answer the question directly; only describe the page content.

# User Query: {text}

# Output only the descriptive paragraph. No introductory text.
# """
#             # Use vLLM for LOCAL HyDE generation
#             if LOCAL:
#                 search_text = ollama_generate_text(
#                 prompt=create_search_prompt,
#                 model="gemma3:1b"
#                 )[0]
#             else:
#                 # Fallback to original DeepInfra if not LOCAL
#                 search_text = DeepInfraInference(
#                     prompt=create_search_prompt,
#                     model_name="Qwen/Qwen3-235B-A22B-Instruct-2507" 
#                 )
            
#             print(f"Search prompt (HyDE via vLLM): {search_text}")
#             # --- END: Query Expansion Logic ---

#             # Encode text
#             # Task 'retrieval.query' optimizes the embedding for finding matching documents
#             embedding = model.encode([search_text], task="retrieval.query")
            
#             # Convert numpy array to list
#             print(f" Generated Jina v4 embedding for text.")
#             return embedding[0].tolist()
        
#         # 4. Handle Image Input (Retrieval Passage)
#         elif image_bytes_list:
#             print(f"Generating Jina v4 embedding (Type: {len(image_bytes_list)} Images)...")
            
#             # Convert raw bytes to PIL Images
#             pil_images = []
#             for img_bytes in image_bytes_list:
#                 pil_images.append(Image.open(io.BytesIO(img_bytes)))
            
#             # Encode images (Batch processing is handled automatically by SentenceTransformer)
#             # Task 'retrieval.passage' optimizes the embedding for being indexed
#             # Note: Ensure the specific Jina model supports image inputs (like Jina-CLIP or specific V4 variants)
#             embeddings = model.encode(pil_images) # task="retrieval.passage" is implied for non-query inputs usually, or add if model supports
            
#             print(f" Generated {len(embeddings)} Jina v4 embeddings for images.")
#             return embeddings.tolist()

#     except Exception as e:
#         print(f"An unexpected error occurred during local embedding generation: {e}")
#         return None


# Configuration for your local vLLM server
VLLM_API_BASE = "http://localhost:4444/v1"  # Adjust port if necessary
VLLM_API_KEY = "EMPTY"  # vLLM usually uses 'EMPTY' or 'token'

def get_image_embedding_jinna_api_local_vllm(
    text: str = None, 
    image_bytes_list: List[bytes] = None, 
    model_name: str = "jinaai/jina-embeddings-v3" # Ensure this model is loaded in vLLM
) -> Union[Optional[List[float]], Optional[List[List[float]]]]:
    """
    Gets a multimodal embedding using a local vLLM server (OpenAI API compatible).
    
    - If 'text' is provided: Generates HyDE description (optional) -> Calls vLLM embedding endpoint.
    - If 'image_bytes_list' is provided: Converts images to Base64 -> Calls vLLM embedding endpoint.

    Args:
        text: The text string to embed.
        image_bytes_list: A list of raw image bytes to embed.
        model_name: The model ID currently served by vLLM.

    Returns:
        A single embedding vector (List[float]) if 'text' was used.
        A list of embedding vectors (List[List[float]]) if 'image_bytes_list' was used.
        None on failure.
    """

    # 1. Input Validation
    if text and image_bytes_list:
        print("Error: Provide either 'text' OR 'image_bytes_list', not both.")
        return None
    if not text and not image_bytes_list:
        print("Error: Must provide either 'text' or 'image_bytes_list'.")
        return None

    # 2. Initialize OpenAI Client pointing to local vLLM
    client = OpenAI(
        base_url=VLLM_API_BASE,
        api_key=VLLM_API_KEY,
    )

    try:
        # 3. Handle Text Input (with HyDE)
        if text:
            print("Processing Text Query with HyDE...")
            
            # --- START: Query Expansion / HyDE Logic ---
            create_search_prompt = f"""
Act as a document search engine. 
Based on the user's query below, generate a detailed paragraph describing the content likely to appear on a document page that answers this query.
User Query: {text}
Output only the descriptive paragraph.
"""
            # NOTE: Ideally, use the 'client.chat.completions' here as well if vLLM handles generation
            # For now, keeping your existing logic structure or assuming a separate generation call exists.
            # If you want to use vLLM for the HyDE generation part too:
            try:
                hyde_response = client.chat.completions.create(
                    model=model_name, # Or a specific chat model loaded in vLLM
                    messages=[{"role": "user", "content": create_search_prompt}],
                    temperature=0.7
                )
                search_text = hyde_response.choices[0].message.content
            except Exception:
                # Fallback if the embedding model doesn't support chat generation
                search_text = text 
                print("HyDE generation failed or model is embedding-only. Using raw text.")

            print(f"Embedding Text (HyDE Result): {search_text[:50]}...")
            # --- END: Query Expansion Logic ---

            # Call vLLM Embeddings Endpoint
            # Note: For Jina v3/v4, you might need to prepend "Query: " depending on how vLLM serves it.
            response = client.embeddings.create(
                input=[search_text],
                model=model_name
            )
            
            embedding = response.data[0].embedding
            print(f"Generated embedding for text.")
            return embedding
        
        # 4. Handle Image Input
        elif image_bytes_list:
            print(f"Processing {len(image_bytes_list)} Images via vLLM API...")
            
            input_payload = []
            
            for img_bytes in image_bytes_list:
                # Convert bytes to Base64 string
                base64_image = base64.b64encode(img_bytes).decode('utf-8')
                
                # Check model requirements: 
                # Some vLLM multimodal forks expect the raw base64 string, 
                # others expect a data URI scheme "data:image/jpeg;base64,..."
                # Here we assume standard Base64 string.
                input_payload.append(base64_image)
            
            # Call vLLM Embeddings Endpoint
            response = client.embeddings.create(
                input=input_payload,
                model=model_name
            )
            
            # Extract embeddings preserving order
            embeddings = [data.embedding for data in response.data]
            
            print(f"Generated {len(embeddings)} embeddings for images.")
            return embeddings

    except Exception as e:
        print(f"An unexpected error occurred during API embedding generation: {e}")
        return None



# ==============================================================================
#  LEGACY & NEW: DATABASE SAVE/SEARCH
# ==============================================================================
    
def encode_text_for_embedding(text: str, target_dimensions: int = 2048, is_query: bool = False) -> list[float]:
    """
    Convert text into an embedding vector using pre-loaded model (FAST).
    Falls back to Ollama if pre-loaded model unavailable.
    
    Args:
        text: ข้อความที่ต้องการ embedding
        target_dimensions: จำนวน dimensions ที่ต้องการ (default: 2048 สำหรับ verified_answers)
        is_query: True = ค้นหา (retrieval.query), False = บันทึกเอกสาร (retrieval.passage)
                  การใช้ task ที่ถูกต้องช่วยให้ cross-lingual search ทำงานได้ดี
    
    Returns:
        list[float]: embedding vector (2048 dims)
    """
    global model  # Use the pre-loaded model
    
    # ตรวจสอบว่า text ว่างหรือไม่
    if not text or not text.strip():
        print(f"❌ ERROR: Empty text provided!")
        raise ValueError("Cannot create embedding from empty text")
    
    # เลือก task - Jina v4 ใช้ 'retrieval' สำหรับทั้ง query และ document
    # Note: is_query ยังคงใช้ประโยชน์สำหรับ logging และ API fallback
    task = 'retrieval'  # Jina v4 ใช้ task เดียวกันสำหรับทั้ง query และ passage
    
    try:
        # Use the pre-loaded Jina embedding model (FAST - no reload)
        if model is not None:
            mode_str = 'QUERY' if is_query else 'DOCUMENT'
            print(f"⚡ Using PRE-LOADED Jina model (task={task}, mode={mode_str}) for embedding...")
            embedding = model.encode(text, task=task)
            embedding_list = embedding.tolist()
            
            # Adjust dimensions to target_dimensions (2048)
            current_dim = len(embedding_list)
            if current_dim < target_dimensions:
                # Pad with zeros
                padding = [0.0] * (target_dimensions - current_dim)
                embedding_list.extend(padding)
                print(f"✅ Padded embedding from {current_dim} to {target_dimensions} dimensions")
            elif current_dim > target_dimensions:
                # Truncate
                embedding_list = embedding_list[:target_dimensions]
                print(f"✅ Truncated embedding from {current_dim} to {target_dimensions} dimensions")
            
            return embedding_list
        else:
            print("⚠️ Model not initialized (model=None). Using Jinna API (Provider API) fallback ...")
            # ส่ง is_query ไปยัง API เพื่อใช้ task ที่ถูกต้อง
            if is_query:
                embedding_list = get_image_embedding_jinna_api(search_text=text)  # retrieval.query
            else:
                embedding_list = get_image_embedding_jinna_api(text=text)  # retrieval.passage
            if embedding_list and len(embedding_list) > 0:
                # Adjust dimensions for API fallback too
                current_dim = len(embedding_list)
                if current_dim < target_dimensions:
                    padding = [0.0] * (target_dimensions - current_dim)
                    embedding_list.extend(padding)
                elif current_dim > target_dimensions:
                    embedding_list = embedding_list[:target_dimensions]
                return embedding_list
            else:
                raise ValueError("Ollama returned empty embedding")
            
    except Exception as e:
        print(f"❌ Jina embedding error: {e}. Trying Ollama fallback...")
        try:
            embedding_list = ollama_embed_text(text=text, model="qwen3-embedding:0.6b")
            if embedding_list and len(embedding_list) > 0:
                embedding_result = embedding_list[0]
                # Adjust dimensions for Ollama fallback
                current_dim = len(embedding_result)
                if current_dim < target_dimensions:
                    padding = [0.0] * (target_dimensions - current_dim)
                    embedding_result.extend(padding)
                    print(f"✅ Padded Ollama embedding from {current_dim} to {target_dimensions} dimensions")
                elif current_dim > target_dimensions:
                    embedding_result = embedding_result[:target_dimensions]
                    print(f"✅ Truncated Ollama embedding from {current_dim} to {target_dimensions} dimensions")
                return embedding_result
            else:
                raise ValueError("Ollama returned empty embedding")
        except Exception as ollama_error:
            print(f"❌ ALL embedding methods failed: {ollama_error}")
            raise ValueError(f"Embedding failed: {ollama_error}")

def clean_text(input_text: str) -> str:
    """
    (Original function, unchanged)
    """
    if input_text is None:
        return ""
    # Remove NUL characters
    cleaned = input_text.replace("\x00", "")
    # Optionally, remove all non-printable control chars except common whitespace (\n, \r, \t)
    cleaned = re.sub(r"[^\x20-\x7E\n\r\t]", "", cleaned)
    return cleaned

def save_vector_to_db(user_id, chat_history_id, uploaded_file_id, file_name, text, embedding, page_number: int = -1):
    """
    Save embedding to the 'document_embeddings' table (Legacy).
    
    UPDATED:
    - Takes 'uploaded_file_id' instead of 'chat_history_id' to link to the file.
    - Takes 'page_number' (defaults to -1).
    """
    # Convert Python list to PostgreSQL vector literal (e.g., '[0.1, 0.2, 0.3]')
    vector_literal = f"[{', '.join(map(str, embedding))}]"
    conn = None
    try:
        # Clean text before inserting
        text = clean_text(text)
        
        print(f"🔍 DEBUG save_vector_to_db:")
        print(f"  - user_id: {user_id}")
        print(f"  - chat_history_id: {chat_history_id}")
        print(f"  - uploaded_file_id: {uploaded_file_id}")
        print(f"  - file_name: {file_name}")
        print(f"  - text length: {len(text)}")
        print(f"  - embedding dims: {len(embedding)}")
        print(f"  - page_number: {page_number}")

        conn = get_db_connection()
        if not conn:
            raise Exception("❌ Could not connect to database")
        
        print(f"✅ Database connection established")
        
        cur = conn.cursor()

        query = """
            INSERT INTO document_embeddings (user_id, chat_history_id, uploaded_file_id, extracted_text, embedding, page_number)
            VALUES (%s, %s, %s, %s, %s, %s)
        """
        
        print(f"🔍 Executing query with parameters: user_id={user_id}, chat_history_id={chat_history_id}, uploaded_file_id={uploaded_file_id}")

        cur.execute(query, (user_id, chat_history_id, uploaded_file_id, text, vector_literal, page_number))

        conn.commit()
        cur.close()
        print("✅ Legacy vector saved to database successfully (page: -1).")
        return True
        
    except Exception as e:
        print(f"❌ Failed to save legacy vector: {e}")
        import traceback
        traceback.print_exc()
        if conn:
            conn.rollback()
        return False
        
    finally:
        if conn:
            conn.close()
            print("✅ Database connection closed")

# --- NEW: Save Page Vector Function ---
def save_page_vector_to_db(user_id, chat_history_id, uploaded_file_id, page_number, embedding):
    """
    Save image embedding to the 'document_page_embeddings' table (New).
    """
    vector_literal = f"[{', '.join(map(str, embedding))}]"
    conn = None
    try:
        conn = get_db_connection()
        if not conn:
            raise Exception("Could not connect to database")
            
        cur = conn.cursor()

        print(f"user_id:{user_id}")
        print(f"chat_id:{chat_history_id}")

        query = """
            INSERT INTO document_page_embeddings (user_id, chat_history_id, uploaded_file_id, page_number, embedding)
            VALUES (%s, %s, %s, %s, %s)
        """

        cur.execute(query, (user_id, chat_history_id, uploaded_file_id, page_number, vector_literal))

        conn.commit()
        cur.close()
        print(f"✅ Page image vector saved to database (Page: {page_number}).")
    except Exception as e:
        print(f"❌ Failed to save page image vector (Page: {page_number}): {e}")
        if conn:
            conn.rollback()
    finally:
        if conn:
            conn.close()

# Search Text (Legacy)
def search_similar_documents_by_chat(query_text: str, user_id: int, chat_history_id: int, top_k: int = 5, threshold_text: float = 0.5):
    """
    Search (Legacy) from 'document_embeddings' table.
    
    UPDATED:
    - Joins with 'uploaded_files' to filter by 'chat_history_id'.
    """
    # Step 1: Encode the query text to a vector
    # query_embedding = encode_text_for_embedding(query_text)
    if not LOCAL:
        if IFXGPT:
            query_embedding = IFXGPTEmbedding(inputs=[query_text])[0]
        else:
            query_embedding = get_image_embedding_jinna_api(search_text=query_text)
    else :
        query_embedding = get_image_embedding_jinna_api_local(search_text=query_text)
    query_vector = f"[{', '.join(map(str, query_embedding))}]"
    
    conn = None
    try:
        conn = get_db_connection()
        if not conn:
            raise Exception("Could not connect to database")
            
        cur = conn.cursor()

        # Step 2: Search within same user and same chat
        # JOIN with uploaded_files to filter by chat_history_id
        print(f"🔍 Searching legacy documents for chat_id={chat_history_id}, threshold={threshold_text}, top_k={top_k}...")
        query = """
            SELECT 
                t1.id AS page_embedding_id, 
                t2.file_name, 
                t2.object_name,
                t1.page_number,
                t1.extracted_text,
                t1.embedding <-> %s AS distance
            FROM document_embeddings AS t1
            INNER JOIN uploaded_files AS t2 ON t1.uploaded_file_id = t2.id
            WHERE t2.user_id = %s
                AND t2.chat_history_id = %s
                AND (t1.embedding <-> %s) <= %s
            ORDER BY distance
            LIMIT %s
        """
        cur.execute(query, (query_vector, user_id, chat_history_id, query_vector, threshold_text, top_k))
        results = cur.fetchall()
        cur.close()

        # Step 3: Return results
        print(f"✅ Found {len(results)} legacy document matches")
        
        # Show search results for debugging
        if results:
            print("Search Results:")
            for row in results:
                print(f"  - {row[1]} (distance: {row[5]:.4f})")

        return [
            {'id': row[0], 'file_name': row[1],'object_name' : row[2],'page_number' : row[3], 'text': row[4], 'distance': row[5]}
            for row in results
        ]

    except Exception as e:
        print("❌ Failed to perform legacy similarity search:", e)
        return []
    finally:
        if conn:
            conn.close()

# --- NEW: Search Similar Pages Function ---
# Search Page (New)
from typing import List, Dict, Any
# Assuming get_image_embedding_jinna_api and get_db_connection are defined elsewhere
# import { get_image_embedding_jinna_api, get_db_connection } from ...

def search_similar_pages(query_text: str, user_id: int, chat_history_id: int, top_k: int = 5, threshold: float = 1.0) -> List[Dict[str, Any]]:
    """
    Search (New) from 'document_page_embeddings' table.
    
    Args:
        query_text: The user's text question.
        user_id: User's ID.
        chat_history_id: The current chat ID.
        top_k: Max number of pages to return *before* normalization/filtering.
        threshold: Max L2 distance for the *initial* SQL query. Results > threshold are excluded by SQL.
    
    Returns:
        A list of dicts, e.g.:
        [{'page_id': 12, 'file_name': 'report.pdf', ..., 'distance': 0.25, 'normalized_distance': 0.1}, ...]
    """
    # Step 1: Encode the query text using the *CLIP* model
    if not LOCAL:
        if IFXGPT: # Embedding
            query_embedding = IFXGPTEmbedding(inputs=[query_text])[0]
        else:
            query_embedding = get_image_embedding_jinna_api(search_text=query_text)
    else :
        query_embedding = get_image_embedding_jinna_api_local(search_text=query_text)
    if not query_embedding:
        print("❌ Failed to get CLIP embedding for query.")
        return []
        
    query_vector = f"[{', '.join(map(str, query_embedding))}]"
    
    conn = None
    try:
        conn = get_db_connection()
        if not conn:
            raise Exception("Could not connect to database")
            
        cur = conn.cursor()

        # Step 2: Search within same user and same chat
        # The SQL query remains the same, using the 'threshold' for a coarse first pass
        # and 'top_k' to limit the initial result set.
        query = """
            SELECT 
                t1.id AS page_embedding_id, 
                t2.file_name, 
                t2.object_name,
                t1.page_number,
                t1.embedding <-> %s AS distance
            FROM document_page_embeddings AS t1
            INNER JOIN uploaded_files AS t2 ON t1.uploaded_file_id = t2.id
            WHERE t2.user_id = %s 
              AND t2.chat_history_id = %s
              AND (t1.embedding <-> %s) <= %s
            ORDER BY distance
            LIMIT %s
        """

        cur.execute(query, (query_vector, user_id, chat_history_id, query_vector, threshold, top_k))
        results = cur.fetchall() # This is the raw list of tuples
        cur.close()
        print(f"✅ Found {len(results)} raw pages within SQL threshold {threshold}.")

        # --- START: New logic for Normalization and 0.3 Threshold ---
        if not results:
            print("ℹ️ No pages found, returning empty list.")
            return []

        # Get all distances to find min/max for normalization
        all_distances = [row[4] for row in results]
        min_dist = min(all_distances)
        max_dist = max(all_distances)
        dist_range = max_dist - min_dist
        
        print(f"ℹ️ Normalizing distances: min={min_dist:.4f}, max={max_dist:.4f}, range={dist_range:.4f}")

        processed_results = []
        for row in results:
            original_distance = row[4]
            
            # Handle division by zero if all distances are identical (dist_range == 0)
            if dist_range == 0:
                # If all distances are the same, all results are equally good
                similarity_score = 1.0
            else:
                # Min-Max normalization: (value - min) / (max - min)
                normalized_distance = (original_distance - min_dist) / dist_range
                # Convert distance to similarity: lower distance = higher similarity
                similarity_score = 1.0 - normalized_distance
            
            processed_results.append({
                'page_embedding_id': row[0],
                'file_name': row[1],
                'object_name': row[2],
                'page_number': row[3],
                'distance': original_distance,
                'similarity_score': similarity_score,
                'double-precision': f'{similarity_score:.10f}'  # Show full precision
            })
        
        print(f"processed_results : {processed_results}")

        # Now, filter based on the *new* similarity threshold (higher is better)
        similarity_threshold = 0.5  # Keep results with similarity >= 0.5
        final_filtered_results = [
            item for item in processed_results 
            if item['similarity_score'] >= similarity_threshold
        ]
        
        print(f"ℹ️ Filtered to {len(final_filtered_results)} pages (similarity_score >= {similarity_threshold}).")
        # --- END: New logic ---

        # Step 3: Return *final filtered* results
        return final_filtered_results

    except Exception as e:
        print("❌ Failed to perform page similarity search:", e)
        return []
    finally:
        if conn:
            conn.close()


# ==============================================================================
#  SEARCH BY ACTIVE USER FUNCTIONS (METHOD: searchDoc)
# ==============================================================================

def search_similar_documents_by_active_user(query_text: str, user_id: int, top_k: int = 5, threshold_text: float = 0.5):
    """
    Legacy Text Search: Finds text chunks in files where the user is an 'active_user'.
    """
    # Encode query

    # Encoding using Qwen3-0.6b-embedding
    # query_embedding = encode_text_for_embedding(text=query_text)

    # Encode using jinna Text-Image-Embedding
    if not LOCAL:
        if IFXGPT:
            query_embedding = IFXGPTEmbedding(inputs=[query_text])[0]
        else:
            query_embedding = get_image_embedding_jinna_api(search_text=query_text)
    else :
        query_embedding = get_image_embedding_jinna_api_local(search_text=query_text)
    query_vector = f"[{', '.join(map(str, query_embedding))}]"  

    conn = None
    try:
        conn = get_db_connection()
        if not conn: raise Exception("DB Connection failed")
        cur = conn.cursor()

        # JOIN uploaded_files and filter by active_users array using ANY()
        query = """
            SELECT 
                t1.id AS page_embedding_id, 
                t2.file_name, 
                t2.object_name,
                t1.page_number,
                t1.extracted_text,
                t1.embedding <-> %s AS distance
            FROM document_embeddings AS t1
            INNER JOIN uploaded_files AS t2 ON t1.uploaded_file_id = t2.id
            WHERE %s = ANY(t2.active_users)
              AND (t1.embedding <-> %s) <= %s
            ORDER BY distance
            LIMIT %s
        """
        cur.execute(query, (query_vector, user_id, query_vector, threshold_text, top_k))
        results = cur.fetchall()
        cur.close()

        # Show search results for debugging
        print("Search Results:")
        for row in results:
            print(row[1] + " ==> " + str(row[5]))

        return [
            {'id': row[0], 'file_name': row[1],'object_name' : row[2],'page_number' : row[3], 'text': row[4], 'distance': row[5]}
            for row in results
        ]
    except Exception as e:
        print(f"Error in search_similar_documents_by_active_user: {e}")
        return []
    finally:
        if conn: conn.close()


def search_similar_pages_by_active_user(query_text: str, user_id: int, top_k: int = 5, threshold: float = 1.0) -> List[Dict[str, Any]]:
    """
    New Page Search: Finds document pages (images) in files where the user is an 'active_user'.
    """
    # 1. Generate Query Embedding (CLIP/Jina)
    if not LOCAL:
        if IFXGPT: # Embedding
            query_embedding = IFXGPTEmbedding(inputs=[query_text])[0]
        else:
            query_embedding = get_image_embedding_jinna_api(search_text=query_text)
    else:
        query_embedding = get_image_embedding_jinna_api_local(search_text=query_text)

    if not query_embedding: return []

    query_vector = f"[{', '.join(map(str, query_embedding))}]"

    conn = None
    try:
        conn = get_db_connection()
        if not conn: raise Exception("DB Connection failed")
        cur = conn.cursor()

        # 2. Search DB (Filter by active_users)
        query = """
            SELECT 
                t1.id AS page_embedding_id, 
                t2.file_name, 
                t2.object_name,
                t1.page_number,
                t1.embedding <-> %s AS distance
            FROM document_page_embeddings AS t1
            INNER JOIN uploaded_files AS t2 ON t1.uploaded_file_id = t2.id
            WHERE %s = ANY(t2.active_users)
              AND (t1.embedding <-> %s) <= %s
            ORDER BY distance
            LIMIT %s
        """
        cur.execute(query, (query_vector, user_id, query_vector, threshold, top_k))
        results = cur.fetchall()
        cur.close()

        # 3. Normalize & Filter (Same logic as standard search)
        if not results: return []

        all_distances = [row[4] for row in results]
        min_dist = min(all_distances)
        max_dist = max(all_distances)
        dist_range = max_dist - min_dist
        
        processed_results = []
        for row in results:
            original_distance = row[4]
            
            if dist_range == 0:
                similarity_score = 1.0
            else:
                normalized_distance = (original_distance - min_dist) / dist_range
                similarity_score = 1.0 - normalized_distance
            
            processed_results.append({
                'page_embedding_id': row[0],
                'file_name': row[1],
                'object_name': row[2],
                'page_number': row[3],
                'distance': original_distance,
                'similarity_score': similarity_score,
                'double-precision': f'{similarity_score:.10f}'
            })

        print(f"processed_results : {processed_results}")
        
        # Filter by similarity threshold (higher is better)
        similarity_threshold = 0.5
        final_filtered_results = [item for item in processed_results if item['similarity_score'] >= similarity_threshold]
        
        return final_filtered_results

    except Exception as e:
        print(f"Error in search_similar_pages_by_active_user: {e}")
        return []
    finally:
        if conn: conn.close()


def search_similar_documents_by_active_user_all(query_text: str, user_id: int, top_k: int = 5, threshold_text: float = 0.5):
    """
    Legacy Text Search: Finds text chunks in all files where the user is an 'active_user'.
    """
    # Encode query
    # query_embedding = encode_text_for_embedding(query_text)
    if not LOCAL:
        if IFXGPT:
            query_embedding = IFXGPTEmbedding(inputs=[query_text])[0]
        else:
            query_embedding = get_image_embedding_jinna_api(search_text=query_text)
    else :
        query_embedding = get_image_embedding_jinna_api_local(search_text=query_text)
    query_vector = f"[{', '.join(map(str, query_embedding))}]"

    conn = None
    try:
        conn = get_db_connection()
        if not conn: raise Exception("DB Connection failed")
        cur = conn.cursor()

        # JOIN uploaded_files and filter by active_users array using ANY()
        query = """
            SELECT 
                t1.id AS page_embedding_id, 
                t2.file_name, 
                t2.object_name,
                t1.page_number,
                t1.extracted_text,
                t1.embedding <-> %s AS distance
            FROM document_embeddings AS t1
            INNER JOIN uploaded_files AS t2 ON t1.uploaded_file_id = t2.id
            WHERE %s = t2.chat_history_id
              AND (t1.embedding <-> %s) <= %s
            ORDER BY distance
            LIMIT %s
        """
        cur.execute(query, (query_vector, -1, query_vector, threshold_text, top_k))
        results = cur.fetchall()
        cur.close()

        # Show search results for debugging
        print("Search Results:")
        for row in results:
            print(row[1] + " ==> " + str(row[5]))

        return [
            {'id': row[0], 'file_name': row[1],'object_name' : row[2],'page_number' : row[3], 'text': row[4], 'distance': row[5]}
            for row in results
        ]
    except Exception as e:
        print(f"Error in search_similar_documents_by_active_user_all: {e}")
        return []
    finally:
        if conn: conn.close()


def search_similar_pages_by_active_user_all(query_text: str, user_id: int, top_k: int = 5, threshold: float = 1.0) -> List[Dict[str, Any]]:
    """
    New Page Search: Finds document pages (images) in all files where the user is an 'active_user'.
    """
    # 1. Generate Query Embedding (CLIP/Jina)
    if not LOCAL:
        if IFXGPT: # Embedding
            query_embedding = IFXGPTEmbedding(inputs=[query_text])[0]
        else:
            query_embedding = get_image_embedding_jinna_api(search_text=query_text)
    else:
        query_embedding = get_image_embedding_jinna_api_local(search_text=query_text)

    if not query_embedding: return []

    query_vector = f"[{', '.join(map(str, query_embedding))}]"

    conn = None
    try:
        conn = get_db_connection()
        if not conn: raise Exception("DB Connection failed")
        cur = conn.cursor()

        # 2. Search DB (Filter by active_users)
        query = """
            SELECT 
                t1.id AS page_embedding_id, 
                t2.file_name, 
                t2.object_name,
                t1.page_number,
                t1.embedding <-> %s AS distance
            FROM document_page_embeddings AS t1
            INNER JOIN uploaded_files AS t2 ON t1.uploaded_file_id = t2.id
            WHERE %s = t2.chat_history_id
              AND (t1.embedding <-> %s) <= %s
            ORDER BY distance
            LIMIT %s
        """
        cur.execute(query, (query_vector, -1, query_vector, threshold, top_k))
        results = cur.fetchall()
        cur.close()

        # 3. Normalize & Filter (Same logic as standard search)
        if not results: return []

        all_distances = [row[4] for row in results]
        min_dist = min(all_distances)
        max_dist = max(all_distances)
        dist_range = max_dist - min_dist
        
        processed_results = []
        for row in results:
            original_distance = row[4]
            
            if dist_range == 0:
                similarity_score = 1.0
            else:
                normalized_distance = (original_distance - min_dist) / dist_range
                similarity_score = 1.0 - normalized_distance
            
            processed_results.append({
                'page_embedding_id': row[0],
                'file_name': row[1],
                'object_name': row[2],
                'page_number': row[3],
                'distance': original_distance,
                'similarity_score': similarity_score,
                'double-precision': f'{similarity_score:.10f}'
            })

        print(f"processed_results : {processed_results}")
        
        # Filter by similarity threshold (higher is better)
        similarity_threshold = 0.5
        final_filtered_results = [item for item in processed_results if item['similarity_score'] >= similarity_threshold]
        return final_filtered_results
    except Exception as e:
        print(f"Error in search_similar_pages_by_active_user_all: {e}")
        return []
    finally:
        if conn: conn.close()

# def search_similar_pages(query_text: str, user_id: int, chat_history_id: int, top_k: int = 5, threshold: float = 1.0) -> List[Dict[str, Any]]:
#     """
#     Search (New) from 'document_page_embeddings' table.
    
#     Args:
#         query_text: The user's text question.
#         user_id: User's ID.
#         chat_history_id: The current chat ID.
#         top_k: Max number of pages to return *before* mean filtering.
#         threshold: Max distance (e.g., 1.0 for L2). Results with distance > threshold are excluded.
    
#     Returns:
#         A list of dicts, e.g.:
#         [{'page_id': 12, 'file_name': 'report.pdf', 'object_name': '...', 'page_number': 5, 'distance': 0.25}, ...]
#     """
#     # Step 1: Encode the query text using the *CLIP* model
#     query_embedding = get_image_embedding_jinna_api(text=query_text)
#     if not query_embedding:
#         print("❌ Failed to get CLIP embedding for query.")
#         return []
        
#     query_vector = f"[{', '.join(map(str, query_embedding))}]"
    
#     conn = None
#     try:
#         conn = get_db_connection()
#         if not conn:
#             raise Exception("Could not connect to database")
            
#         cur = conn.cursor()

#         # Step 2: Search within same user and same chat
#         # JOIN with uploaded_files to filter by chat_history_id and get object_name
#         # Use <-> (L2 distance) operator.
#         query = """
#             SELECT 
#                 t1.id AS page_embedding_id, 
#                 t2.file_name, 
#                 t2.object_name,
#                 t1.page_number,
#                 t1.embedding <-> %s AS distance
#             FROM document_page_embeddings AS t1
#             INNER JOIN uploaded_files AS t2 ON t1.uploaded_file_id = t2.id
#             WHERE t2.user_id = %s 
#               AND t2.chat_history_id = %s
#               AND (t1.embedding <-> %s) <= %s
#             ORDER BY distance
#             LIMIT %s
#         """

#         cur.execute(query, (query_vector, user_id, chat_history_id, query_vector, threshold, top_k))
#         results = cur.fetchall() # This is the raw list of tuples
#         cur.close()
#         print(f"✅ Found {len(results)} similar pages within threshold {threshold}.")
#         print("Raw Results:", results)

#         # --- START: New logic to filter by mean distance ---
#         if not results:
#             print("ℹ️ No pages found, returning empty list.")
#             return []

#         # Calculate mean distance from the initial results
#         all_distances = [row[4] for row in results]
#         mean_distance = sum(all_distances) / len(all_distances)
#         print(f"ℹ️ Calculated mean distance: {mean_distance}")

#         # Filter the results list based on the mean
#         filtered_results = [row for row in results if row[4] <= mean_distance]
#         print(f"ℹ️ Filtered to {len(filtered_results)} pages (distance <= mean).")
#         # --- END: New logic ---

#         # Step 3: Return *filtered* results
#         return [
#             {
#                 'page_embedding_id': row[0],
#                 'file_name': row[1],
#                 'object_name': row[2],
#                 'page_number': row[3], # 1-indexed
#                 'distance': row[4]
#             }
#             for row in filtered_results # Iterate over the new filtered list
#         ]

#     except Exception as e:
#         print("❌ Failed to perform page similarity search:", e)
#         return []
#     finally:
#         if conn:
#             conn.close()


# def search_similar_pages(query_text: str, user_id: int, chat_history_id: int, top_k: int = 5, threshold: float = 1.0) -> List[Dict[str, Any]]:
#     """
#     Search (New) from 'document_page_embeddings' table.
    
#     Args:
#         query_text: The user's text question.
#         user_id: User's ID.
#         chat_history_id: The current chat ID.
#         top_k: Max number of pages to return.
#         threshold: Max distance (e.g., 1.0 for L2). Results with distance > threshold are excluded.
    
#     Returns:
#         A list of dicts, e.g.:
#         [{'page_id': 12, 'file_name': 'report.pdf', 'object_name': '...', 'page_number': 5, 'distance': 0.25}, ...]
#     """
#     # Step 1: Encode the query text using the *CLIP* model
#     query_embedding = get_image_embedding_jinna_api(text=query_text)
#     if not query_embedding:
#         print("❌ Failed to get CLIP embedding for query.")
#         return []
        
#     query_vector = f"[{', '.join(map(str, query_embedding))}]"
    
#     conn = None
#     try:
#         conn = get_db_connection()
#         if not conn:
#             raise Exception("Could not connect to database")
            
#         cur = conn.cursor()

#         # Step 2: Search within same user and same chat
#         # JOIN with uploaded_files to filter by chat_history_id and get object_name
#         # Use <-> (L2 distance) operator.
#         query = """
#             SELECT 
#                 t1.id AS page_embedding_id, 
#                 t2.file_name, 
#                 t2.object_name,
#                 t1.page_number,
#                 t1.embedding <-> %s AS distance
#             FROM document_page_embeddings AS t1
#             INNER JOIN uploaded_files AS t2 ON t1.uploaded_file_id = t2.id
#             WHERE t2.user_id = %s 
#               AND t2.chat_history_id = %s
#               AND (t1.embedding <-> %s) <= %s
#             ORDER BY distance
#             LIMIT %s
#         """

#         cur.execute(query, (query_vector, user_id, chat_history_id, query_vector, threshold, top_k))
#         results = cur.fetchall()
#         cur.close()
#         print(f"✅ Found {len(results)} similar pages within threshold {threshold}.")
#         print("Results:", results)

#         # Step 3: Return results
#         return [
#             {
#                 'page_embedding_id': row[0],
#                 'file_name': row[1],
#                 'object_name': row[2],
#                 'page_number': row[3], # 1-indexed
#                 'distance': row[4]
#             }
#             for row in results
#         ]

#     except Exception as e:
#         print("❌ Failed to perform page similarity search:", e)
#         return []
#     finally:
#         if conn:
#             conn.close()


# --- NEW: VLM Page Processing Orchestrator ---
def process_pages_with_vlm(search_results: List[Dict[str, Any]], original_query: str) -> str:
    """
    Orchestrates the new RAG flow:
    1. Takes search results (list of pages).
    2. Fetches the original file for each from MinIO.
    3. Extracts the specific page as an image.
    4. Sends all images + query to a VLM.
    5. Returns the VLM's text response.
    """
    if not search_results:
        return "I found no relevant document pages for your query."

    print(f"🚀 Processing {len(search_results)} relevant pages with VLM...")
    
    # Use a dictionary to cache file fetches from MinIO
    file_cache = {} 
    image_bytes_list = []
    
    page_references = [] # To tell the VLM what it's looking at

    for result in search_results:
        object_name = result['object_name']
        page_num_1_idx = result['page_number']
        file_name = result['file_name']
        
        # 1. Fetch file from cache or MinIO
        if object_name not in file_cache:
            print(f"Fetching '{object_name}' from MinIO...")
            file_bytes = get_file_from_minio(object_name)
            if file_bytes:
                file_cache[object_name] = file_bytes
            else:
                print(f"⚠️ Could not fetch file {object_name}. Skipping page {page_num_1_idx}.")
                continue
        
        file_bytes = file_cache[object_name]
        
        # 2. Extract the specific page as an image
        # Page number is 1-indexed in DB, convert to 0-indexed for fitz
        page_num_0_idx = page_num_1_idx - 1 
        image_bytes = convert_pdf_page_to_image(file_bytes, page_num_0_idx, dpi=100)
        
        if image_bytes:
            image_bytes_list.append(image_bytes)
            page_references.append(f"- '{file_name}' (Page {page_num_1_idx})")
        else:
            print(f"⚠️ Failed to render page {page_num_1_idx} from {file_name}. Skipping.")
            
    if not image_bytes_list:
        return "I found relevant pages but could not render them as images for analysis."

    # 3. Send all images to the VLM
    print(f"Sending {len(image_bytes_list)} images to VLM...")

#     system_prompt = (
# """You are an expert Document Analyst AI. Your mission is to meticulously analyze an image of a document and convert its **entire content** into a structured Markdown format. You must process the document sequentially from top to bottom, preserving the original order of all elements.

# **Your primary directive is to perform a literal, verbatim extraction. You must not summarize, interpret, rephrase, or omit any content.**

# Your response must be a single, complete Markdown document representing the source file.

# ### Your Guiding Principles

# 1.  **Absolute Sequential Order:** Process the document from its absolute top to its absolute bottom, transcribing elements in the exact order they appear. If content is arranged visually (e.g., side-by-side columns), you must transcribe the content in a logical reading order (e.g., top-to-bottom of the left column, then top-to-bottom of the right column). The order of elements in your output **must** perfectly match the original.

# 2.  **Verbatim Text Transcription:** All text elements must be transcribed **verbatim (exactly as written)** and converted into the appropriate Markdown format. This includes:

#       * Headings (e.g., `# Title`, `## Subtitle`)
#       * Paragraphs (plain text)
#       * Lists (bulleted `*` or numbered `1.`)
#       * Code blocks (using \`\`\` fences)
#       * Bold (`**text**`) and Italic (`*text*`) styling.
#       * Any and all text labels, captions, or annotations, in the exact place they appear.

# 3.  **Literal Deconstruction of Non-Textual Elements (CRITICAL):** This is **not a summary**. This is a **textual conversion** of visual information. For any visual element that cannot be represented as text (such as images, photographs, charts, **circuit diagrams, schematics,** signatures, stamps, or logos), you must:

#       * Provide a **hyper-detailed, component-by-component, literal transcription** of the element's content.
#       * **Zero Summarization:** Your transcription must be exhaustive and deconstruct the element for someone who cannot see it. You must transcribe all labels, all data points, and all connections.
#       * **For Diagrams/Schematics:** This is the most critical task. You must *trace* and *transcribe* **every single connection** between all labeled components. This is a textual representation of the visual data, not a summary of its purpose.
#           * List each component by its label (e.g., "Arduino Pin A0/14/PC0", "Resistor R1", "Dip Switch 1").
#           * Explicitly trace what each pin connects to, including any intermediate components. (e.g., "Arduino Pin 0/PD0 connects to one side of the first switch in 'Dip Switch 1' AND to one end of a '10K-ohm' resistor. The other end of this resistor connects to '+5V'. The other side of the first switch connects to 'GND'.").
#       * **For Charts/Graphs:** State the chart type (bar, line, pie). Transcribe all axes labels, the title, and the legend. Then, list the data points or relationships shown, one by one.
#       * Enclose this entire literal deconstruction within the specific tags defined in Principle 4.

# 4.  **Tagging Format:** You must use the following tags to enclose your non-textual deconstructions. **The text *inside* the tags is the literal transcription of the visual, not a summary.**

#       * **Tables:** Enclose the entire Markdown table within `<table>` and `</table>`.
#         ```markdown
#         <table>
#         | Header 1 | Header 2 |
#         |---|---|
#         | Data 1 | Data 2 |
#         </table>
#         ```
#       * **Charts:**
#         ```markdown
#         <chart>A vertical bar chart titled 'Sales'. The X-axis is 'Month' (Jan, Feb, Mar). The Y-axis is 'Revenue'. The bar for Feb is taller than Jan.</chart>
#         ```
#       * **Images/Photos:**
#         ```markdown
#         <image>A photograph of a white coffee mug with a blue logo, sitting on a wooden desk next to a laptop.</image>
#         ```
#       * **Diagrams/Schematics (Your Most Important Tag):**
#         ```markdown
#         <diagram>
#         A hyper-detailed, connection-by-connection transcription of the schematic.
#         **Component 1 (e.g., Arduino):**
#         * Pin 1 connects to...
#         * Pin 2 connects to Resistor R1...
#         **Component 2 (e.g., 7-Segment Display):**
#         * Pin 'A' connects to the other end of Resistor R1...
#         * Pin 'Common' connects to GND...
#         (This transcription must trace all connections literally from the image.)
#         </diagram>
#         ```
#       * **Logos:**
#         ```markdown
#         <logo>A circular company logo with the text 'Innovate Corp' and a gear icon in the center.</logo>
#         ```
#       * **Signatures:**
#         ```markdown
#         <signature>A handwritten, illegible signature in blue ink.</signature>
#         ```
#       * **Stamps:**
#         ```markdown
#         <stamp>A red circular stamp with the text 'APPROVED' in the center.</stamp>
#         ```

# ### Relevance Filter
# If a specific user question is provided and this page contains no relevant information to answer it, do not extract content. Instead, return only: "no related data"

# -----

# ### *** Note ***
# You must extract all content from the document.

# ### Gold-Standard Example (Your Target Quality)

# The example you provided is the **perfect** model of this non-summary approach. It correctly provides a literal, connection-by-connection transcription of the visual diagram *first*, and *then* provides a verbatim transcription of all the text that follows it, in the correct order.

# **Your required output must be exactly in this format and at this level of detail:**

# ```markdown
# <diagram>
# A detailed circuit schematic showing connections between an Arduino-style board, a common cathode 7-Segment display, and an 8-position Dip Switch block labeled "ดิปสวิตช์1".

# **Connections from Arduino Board:**

# * **To 7-Segment Display (via 300-ohm resistors):**
#     * Pin A0/14/PC0 connects to a 300-ohm resistor, which connects to the 'A' segment pin.
#     * Pin A1/15/PC1 connects to a 300-ohm resistor, which connects to the 'B' segment pin.
#     * Pin A2/16/PC2 connects to a 300-ohm resistor, which connects to the 'C' segment pin.
#     * Pin A3/17/PC3 connects to a 300-ohm resistor, which connects to the 'D' segment pin.
#     * Pin A4/18/PC4 connects to a 300-ohm resistor, which connects to the 'E' segment pin.
#     * Pin A5/19/PC5 connects to a 300-ohm resistor, which connects to the 'F' segment pin.
#     * Pin 8/PB0 connects to a 300-ohm resistor, which connects to the 'G' segment pin.
# * **To Dip Switch 1 (Pull-Up Configuration):**
#     * Pin 0/PD0 connects to a node. This node is connected to one side of the first switch (leftmost) and also to one end of a 10K-ohm resistor.
#     * Pin 1/PD1 connects to a node. This node is connected to one side of the second switch and also to one end of a 10K-ohm resistor.
#     * Pin 2/PD2 connects to a node. This node is connected to one side of the third switch and also to one end of a 10K-ohm resistor.
#     * Pin 3/PD3 connects to a node. This node is connected to one side of the fourth switch and also to one end of a 10K-ohm resistor.
#     * Pin 4/PD4 connects to a node. This node is connected to one side of the fifth switch and also to one end of a 10K-ohm resistor.
#     * Pin 5/PD5 connects to a node. This node is connected to one side of the sixth switch and also to one end of a 10K-ohm resistor.
#     * Pin 6/PD6 connects to a node. This node is connected to one side of the seventh switch and also to one end of a 10K-ohm resistor.
#     * Pin 7/PD7 connects to a node. This node is connected to one side of the eighth switch (rightmost) and also to one end of a 10K-ohm resistor.

# **Component Connections:**

# * **7-Segment Display:**
#     * Segments A, B, C, D, E, F, and G are connected as described above.
#     * The common cathode pin (at the bottom) is connected to GND.
# * **Dip Switch 1 ("ดิปสวิตช์1"):**
#     * This is an 8-position switch block.
#     * One side of each of the 8 switches is connected to its respective Arduino pin (PD0-PD7) as described above.
#     * The other side of all 8 switches is connected to a common GND line.
# * **Pull-Up Resistors (10K-ohm x 8):**
#     * There are 8 10K-ohm resistors.
#     * One end of each resistor is connected to an Arduino pin node (PD0-PD7).
#     * The other end of all 8 resistors is connected to a common +5V line.

# **Power Pins:**
# * The Arduino board shows a "USB JACK".
# * A power pin block shows 3.3V, 5V, GND, GND, VIN.
# * Unconnected pins on the Arduino board include: SCL, SDA, AREF, GND, 13/PB5, 12/PB4, 11/PB3, 10/PB2, 9/PB1.
# </diagram>

# รูปที่ 1

# Lab 1
# Dip Switch and 7-Segment

# อุปกรณ์
# 1. Arduino Board
# 2. Digital Experiment Board
# 3. 7-Segment Board

# Checkpoint# 1: อ่านค่าสถานะจาก Dip-switch เพื่อแสดงค่าออกที่ 7-Segment
# 1.1 ต่อวงจร ตามรูปที่ 1
# 1.2 ใช้โปรแกรม Arduino IDE เพื่อป้อน Code ด้านล่าง และแก้ไข เพิ่มเติมให้เหมาะสม เพื่อให้ Code สามารถอ่านค่าสถานะตรรกะจาก Dip-Switch ทั้ง 8 แล้วทำการนับจำนวน Switch ที่มี ตรรกะ “HIGH” แล้วแสดงค่าจำนวนนั้นออกสู่ 7-Segment ได้อย่างถูกต้อง
# ```""")
    system_prompt = (
"""You are an expert Document Analyst AI. Your task is to convert document images into structured Markdown sequentially (top-to-bottom).

**CORE DIRECTIVE: Literal, verbatim extraction only. DO NOT summarize, interpret, or omit any content.**

### Operational Rules
1.  **Sequential Order:** Transcribe elements exactly as they appear (top-to-bottom, left-to-right).
2.  **Verbatim Text:** Extract all text exactly as written, preserving Markdown formatting (Headers, Lists, Code blocks, **Bold**, *Italic*).
3.  **Visual Deconstruction (CRITICAL):**
    *   Convert visual elements (charts, diagrams) into text descriptions within specific tags.
    *   **For Schematics/Diagrams:** You must perform a **component-by-component connection trace**. Explicitly state every pin and wire connection (e.g., "Pin A connects to Resistor R1, which connects to GND"). Never summarize the diagram's purpose; describe its structure.

### Required Tags
Use these tags to enclose literal descriptions of non-text elements:
*   `<diagram>`: Detailed schematic connection tracing.
*   `<table>`: Markdown tables.
*   `<chart>`: Type, axes, legend, and data points.
*   `<image>`, `<logo>`, `<signature>`, `<stamp>`: Literal visual description.

### Relevance Filter
If a specific user question is provided and this page contains no relevant information to answer it, do not extract content. Instead, return only: "no related data"

### Gold-Standard Output Example
**Output must match this level of detail:**

```markdown
<diagram>
A detailed circuit schematic showing connections between an Arduino-style board, a 7-Segment display, and a Dip Switch block.

**Connections:**
* **Arduino to 7-Segment:**
    * Pin A0 connects to a 300-ohm resistor -> 'A' segment pin.
    * Pin A1 connects to a 300-ohm resistor -> 'B' segment pin.
    * [List continues for all pins...]
* **Arduino to Dip Switch (Pull-Up):**
    * Pin 0/PD0 connects to Switch 1 and a 10K resistor to +5V.
    * Pin 1/PD1 connects to Switch 2 and a 10K resistor to +5V.
    * [List continues for all pins...]

**Components:**
* **Dip Switch:** 8-position block. Common side connects to GND.
* **Resistors:** 8x 10K-ohm (Pull-up) and 7x 300-ohm (Current limiting).
</diagram>

รูปที่ 1

Lab 1
Dip Switch and 7-Segment

อุปกรณ์
1. Arduino Board
2. Digital Experiment Board
3. 7-Segment Board

Checkpoint# 1: อ่านค่าสถานะจาก Dip-switch เพื่อแสดงค่าออกที่ 7-Segment
1.1 ต่อวงจร ตามรูปที่ 1
1.2 ใช้โปรแกรม Arduino IDE เพื่อป้อน Code ด้านล่าง...
```""")
    
    system_prompt = (
"""You are an expert Document Analyst AI converting images to structured Markdown.

**CORE DIRECTIVE: Extract content sequentially and verbatim. NO summarization, interpretation, or omission.**

### Operational Rules
1.  **Sequential Order:** Transcribe elements top-to-bottom, left-to-right.
2.  **Text Transcription:** Extract text exactly as written, preserving Markdown formatting (Headers, Lists, Code blocks, **Bold**, *Italic*).
3.  **Visual Deconstruction (CRITICAL):**
    *   Convert visuals into literal text descriptions inside specific tags.
    *   **For Schematics/Diagrams:** You must provide a **hyper-detailed, pin-by-pin connection trace**. Explicitly state every wire connection (e.g., "Pin A connects to Resistor R1, which connects to GND"). Describe the structure, not the function.

### Required Tags
*   `<diagram>`: Detailed schematic connection tracing.
*   `<table>`: Markdown tables.
*   `<chart>`: Type, axes, legend, and data points.
*   `<image>` / `<logo>` / `<signature>` / `<stamp>`: Literal visual description.

### Relevance Filter
If a specific user question is provided and this page contains no relevant information to answer it, do not extract content of that section.
""") #**********************************************
    
#     system_prompt = (
# """You are a Document Analyst AI converting images to Markdown.

# **DIRECTIVE: Extract all content sequentially and verbatim. DO NOT summarize.**

# 1. **Text:** Transcribe exactly as written, preserving Markdown formatting (Headers, Lists, Code).
# 2. **Visuals:** Describe structure literally within tags: `<table>`, `<chart>`, `<image>`.
# 3. **Diagrams (CRITICAL):** Inside `<diagram>`, you must trace every connection pin-by-pin (e.g., "Pin A connects to Resistor R1..."). Describe the wiring details, not the concept.

# ### Relevance Filter
# If a specific user question is provided and this page contains no relevant information to answer it, do not extract content. Instead, return only: "no related data"
# """)
    
    reference_text = "\n".join(page_references)
    prompt = f"""
file name : {file_name}
Based on the following {len(image_bytes_list)} document pages:
{reference_text}

Please extract data to markdown (keep all original data ignore not match content) for use to answer the question:
"{original_query}"
*** Do not answer the question. ***
*** Extract all content from the document. ***

### Relevance Filter
*** If a specific user question is provided and this page contains no relevant information to answer it, do not extract content of that section. ***

Provide your response in Markdown format, following the tagging guidelines provided in the system prompt.
"""

    # Use OpenRouterInference, which now accepts a list of images
    if not LOCAL:
        if IFXGPT:
            vlm_response = IFXGPTInference(
            # vlm_response = DeepInfraInference(
                prompt=prompt,
                system_prompt=system_prompt,
                image_bytes_list=image_bytes_list,
                model_name= 'gpt-5-mini'#'Qwen/Qwen3-VL-8B-Instruct'#'qwen/qwen3-vl-8b-instruct'#'Qwen/Qwen2.5-VL-32B-Instruct'#'deepseek-ai/DeepSeek-OCR'#'Qwen/Qwen3-VL-30B-A3B-Instruct'#'deepseek-ai/DeepSeek-V3.2'#'Qwen/Qwen3-VL-30B-A3B-Instruct'#"Qwen/Qwen2.5-VL-32B-Instruct" #'x-ai/grok-4-fast'#"Qwen/Qwen2.5-VL-32B-Instruct" # Use a strong VLM
            )
            print("IFXGPT VLM response received.")
        else:
            # system_prompt = ("You're an image expert."
            #                  "If the image contains text, extract and summarize it...")
            # # Prompt for OpenRouter VLM
            # prompt = ("Please describe the image in detail in a text format that allows you to understand its details.")
            vlm_response = OpenRouterInference(
            # vlm_response = DeepInfraInference(
                prompt=prompt,
                system_prompt=system_prompt,
                image_bytes_list=image_bytes_list,
                model_name= 'google/gemini-2.5-flash-lite'#'Qwen/Qwen3-VL-8B-Instruct'#'qwen/qwen3-vl-8b-instruct'#'Qwen/Qwen2.5-VL-32B-Instruct'#'deepseek-ai/DeepSeek-OCR'#'Qwen/Qwen3-VL-30B-A3B-Instruct'#'deepseek-ai/DeepSeek-V3.2'#'Qwen/Qwen3-VL-30B-A3B-Instruct'#"Qwen/Qwen2.5-VL-32B-Instruct" #'x-ai/grok-4-fast'#"Qwen/Qwen2.5-VL-32B-Instruct" # Use a strong VLM
            )
            print("DeepInfra VLM response received.")
        print(vlm_response)
    else:
        #  # System prompt for OpenRouter VLM
        # system_prompt = ("You're an image expert."
        #                  "If the image contains text, extract and summarize it...")
        # # Prompt for OpenRouter VLM
        # prompt = ("Please describe the image in detail in a text format that allows you to understand its details.")
        vlm_response = ollama_describe_image(
            image_bytes=image_bytes_list,
            model="qwen3-vl:2b-instruct",
            prompt=prompt,
            system_prompt=system_prompt,
        )
        vlm_response = "\n\n".join(vlm_response)
    # for batch processing

    # response_vlm_list = []
    # for idx, resp in enumerate(vlm_response):
    #     response_vlm_list.append({
    #         'id': search_results[idx].get('page_embedding_id'),
    #         'file_name': search_results[idx].get('file_name'),
    #         # 'object_name': search_results[idx].get('object_name'),
    #         # 'page_number': search_results[idx].get('page_number'),
    #         'text': resp,
    #         'distance': search_results[idx].get('distance'),
    #         # 'normalized_distance': search_results[idx].get('normalized_distance'),
    #         # 'description': resp
    #     })
    # vlm_response = "\n\n".join(vlm_response)
    print("✅ VLM processing complete.")
    
    # return vlm_response
    return vlm_response



# ==============================================================================
#  NEW: OLLAMA INFERENCE FUNCTIONS
# ==============================================================================


def ollama_generate_text(prompt: Union[str, List[str]], model: str = "llama3.2:3b", system_prompt: str = "") -> Union[str, List[str]]:
    """
    Generate text using Ollama's API. Supports single prompt or list of prompts.
    
    Args:
        prompt: The user prompt(s) (str or list[str]).
        model: The Ollama model name (e.g., "llama3.2:3b").
        system_prompt: Optional system prompt.
    
    Returns:
        Generated text(s) or error message(s).
    """
    if isinstance(prompt, str):
        prompts = [prompt]
        single = True
    else:
        prompts = prompt
        single = False
    
    results = []
    for p in prompts:
        url = API_OLLAMA
        payload = {
            "model": model,
            "prompt": p,
            "stream": False
        }
        if system_prompt:
            payload["system"] = system_prompt
        
        try:
            response = requests.post(url, json=payload)
            response.raise_for_status()
            data = response.json()
            results.append(data.get("response", "No response generated."))
        except requests.exceptions.RequestException as e:
            results.append(f"Error calling Ollama API: {e}")
        except Exception as e:
            results.append(f"An unexpected error occurred: {e}")
    
    return results[0] if single else results

def ollama_embed_text(text: Union[str, List[str]], model: str = "nomic-embed-text") -> List[List[float]]:
    """
    Generate embeddings for text using Ollama's API. Supports single text or list of texts.
    
    Args:
        text: The text(s) to embed (str or list[str]).
        model: The Ollama embedding model name (e.g., "nomic-embed-text").
    
    Returns:
        List of embeddings (list[list[float]]).
    """
    if isinstance(text, str):
        inputs = [text]
    else:
        inputs = text
    
    url = "http://localhost:11434/api/embed"
    payload = {
        "model": model,
        "input": inputs
    }
    
    try:
        response = requests.post(url, json=payload)
        response.raise_for_status()
        data = response.json()
        embeddings = data.get("embeddings", [])
        return embeddings if embeddings else [[] for _ in inputs]
    except requests.exceptions.RequestException as e:
        print(f"Error calling Ollama embed API: {e}")
        return [[] for _ in inputs]
    except Exception as e:
        print(f"An unexpected error occurred during embedding: {e}")
        return [[] for _ in inputs]

def ollama_describe_image(image_bytes: Union[bytes, List[bytes]], model: str = "llava", prompt: str = "Describe this image in detail.",system_prompt="") -> Union[str, List[str]]:
    """
    Describe an image using Ollama's vision model. Supports single image or list of images.
    
    Args:
        image_bytes: Raw image bytes (bytes or list[bytes]).
        model: The Ollama vision model name (e.g., "llava").
        prompt: The prompt for description.
    
    Returns:
        Description text(s) or error message(s).
    """
    clear_gpu()
    if isinstance(image_bytes, bytes):
        images = [image_bytes]
        single = True
    else:
        images = image_bytes
        single = False
    
    results = []
    for img_bytes in images:
        import base64
        base64_image = base64.b64encode(img_bytes).decode('utf-8')
        
        url = API_OLLAMA
        payload = {
            "model": model,
            "prompt": prompt,
            "system" :system_prompt,
            "images": [base64_image],
            "stream": False
        }
        
        try:
            response = requests.post(url, json=payload)
            response.raise_for_status()
            data = response.json()
            # print("Raw output")
            # print(data)
            results.append(data.get("response", "No description generated."))
        except requests.exceptions.RequestException as e:
            results.append(f"Error calling Ollama vision API: {e}")
        except Exception as e:
            results.append(f"An unexpected error occurred: {e}")
        print(f"Description: {results[0]}")
    clear_gpu()
    return results[0] if single else results

def ollama_embed_image(image_bytes: Union[bytes, List[bytes]], vision_model: str = "llava", embed_model: str = "nomic-embed-text", prompt: str = "Describe this image in detail.") -> List[List[float]]:
    """
    Generate embeddings for an image by first describing it with a vision model, then embedding the description.
    Supports single image or list of images.
    
    Args:
        image_bytes: Raw image bytes (bytes or list[bytes]).
        vision_model: The Ollama vision model for description.
        embed_model: The Ollama embedding model for text.
        prompt: The prompt for image description.
    
    Returns:
        List of embeddings (list[list[float]]).
    """
     # System prompt for OpenRouter VLM
    system_prompt = ("You're an image expert."
                     "If the image contains text, extract and summarize it...")
    # Prompt for OpenRouter VLM
    prompt = ("Please describe the image in detail in a text format that allows you to understand its details.")
    descriptions = ollama_describe_image(image_bytes, vision_model, prompt)
    if isinstance(descriptions, str):
        descriptions = [descriptions]
    
    # Filter out error descriptions
    valid_descriptions = [d for d in descriptions if not d.startswith("Error")]
    if not valid_descriptions:
        return [[] for _ in descriptions]
    
    embeddings = ollama_embed_text(valid_descriptions, embed_model)
    
    # Map back to original order, with empty lists for failed descriptions
    result = []
    idx = 0
    for d in descriptions:
        if d.startswith("Error"):
            result.append([])
        else:
            result.append(embeddings[idx])
            idx += 1
    return result



def IFXGPTInference(prompt: str, system_prompt: str = "", image_bytes_list: List[bytes] = None, model_name: str = "gpt-4o") -> str:
    """
    Perform inference using the Infineon IFX GPT (OpenAI Client).
    """
    messages = []
    if system_prompt:
        messages.append({"role": "system", "content": system_prompt})

    user_content = [{"type": "text", "text": prompt}]

    if image_bytes_list:
        for image_bytes in image_bytes_list:
            try:
                # Guess mime type
                mime_type, _ = mimetypes.guess_type("image.png") # Default
                img = Image.open(io.BytesIO(image_bytes))
                if img.format:
                    mime_type = f"image/{img.format.lower()}"
                
                base64_image = base64.b64encode(image_bytes).decode('utf-8')
                data_url = f"data:{mime_type};base64,{base64_image}"
                
                user_content.append({
                    "type": "image_url",
                    "image_url": {"url": data_url}
                })
            except Exception as e:
                print(f"Error encoding image for IFXGPT: {e}")

    messages.append({"role": "user", "content": user_content})

    try:
        response = client.chat.completions.create(
            model=model_name,
            messages=messages,
            temperature=1.0
        )
        print("Text Extract: ", response.choices[0].message.content)
        return response.choices[0].message.content
    except Exception as e:
        return f"Error calling IFX GPT: {e}"


def IFXGPTEmbedding(
    inputs: List[str],
    model_name: str = "text-embedding-3-small",
) -> List[List[float]]:
    print("string input: ", inputs)
    # 1) validate inputs
    if not inputs:
        raise ValueError("IFXGPTEmbedding: 'inputs' is empty")

    # 2) sanitize (common cause: empty/whitespace chunks)
    clean_inputs = [s.strip() for s in inputs if isinstance(s, str) and s.strip()]
    if not clean_inputs:
        raise ValueError("IFXGPTEmbedding: all inputs are empty/whitespace")
    print("clean input: ", clean_inputs)
    try:
        response = client.embeddings.create(
            model=model_name,
            input=clean_inputs
        )

        # 3) validate response
        data = getattr(response, "data", None)
        if not data:
            # try to show helpful debug info
            raise RuntimeError(f"No embedding data received. Raw response: {response}")

        embeddings = []
        for i, item in enumerate(data):
            emb = getattr(item, "embedding", None)
            if emb is None:
                raise RuntimeError(f"Missing embedding at index {i}. Item: {item}")
            embeddings.append(emb)

        # 4) ensure count matches
        if len(embeddings) != len(clean_inputs):
            raise RuntimeError(
                f"Embedding count mismatch: got {len(embeddings)} for {len(clean_inputs)} inputs"
            )
        return embeddings

    except Exception as e:
        # re-raise so caller can handle it; printing hides failures and causes downstream index errors
        raise RuntimeError(f"IFX GPT Embeddings failed: {e}") from e


# ==============================================================================
#  LEGACY: FILE SYSTEM (Unchanged)
# ==============================================================================

class EditedFileSystem:
    def __init__(self):
        # Define a base directory for file operations
        # You might want to make this configurable
        self.base_dir = "managed_files"
        if not os.path.exists(self.base_dir):
            os.makedirs(self.base_dir)

    def _get_full_path(self, file_name: str) -> str:
        """Helper to get the full path of a file within the managed directory."""
        return os.path.join(self.base_dir, file_name)

    # --- Read file operations ---
    def read_line(self, file_name: str, start_line: int, end_line: int) -> tuple[list, str]:
        """
        Reads specific lines from a file, prepending line numbers to each.
        """
        full_path = self._get_full_path(file_name)
        lines = []
        try:
            with open(full_path, 'r') as f:
                for i, line in enumerate(f, 1):
                    if start_line <= i <= end_line:
                        # clean_line = line.rstrip('\n')
                        lines.append(f"{i}: {line.rstrip()}")
                    elif i > end_line:
                        break # Optimization: stop reading once past the end_line
            return lines, ""
        except FileNotFoundError:
            return [], f"File '{file_name}' not found."
        except Exception as e:
            return [], f"An error occurred while reading file: {e}"

    def read_all(self, file_name: str) -> tuple[list, str]:
        """
        Reads all lines from a file, prepending line numbers to each.
        """
        full_path = self._get_full_path(file_name)
        lines = []
        try:
            with open(full_path, 'r') as f:
                # Enumerate to get line numbers and prepend them
                lines = [f"{i}: {line.rstrip()}" for i, line in enumerate(f, 1)]
            return lines, ""
        except FileNotFoundError:
            return [], f"File '{file_name}' not found."
        except Exception as e:
            return [], f"An error occurred while reading file: {e}"

    def read_start_until_line_n(self, file_name: str, end_line: int) -> tuple[list, str]:
        """
        Reads from the start of the file up to a specified line number,
        prepending line numbers to each.
        """
        # This method already calls read_line, which now handles line numbering
        return self.read_line(file_name, 1, end_line)

    # --- Edit file operations ---
    def edit_line(self, file_name: str, text: str, start_line: int, end_line: int) -> str:
        """
        Edits specific lines in a file by replacing them with the provided text.
        """
        full_path = self._get_full_path(file_name)
        try:
            with open(full_path, 'r') as f:
                lines = f.readlines()

            start_idx = start_line - 1
            end_idx = end_line - 1

            if not (0 <= start_idx <= end_idx < len(lines)):
                return f"Line range {start_line}-{end_line} is out of bounds for file '{file_name}'."

            new_lines_to_insert = [line + '\n' for line in text.splitlines()]
            if not new_lines_to_insert:
                new_lines_to_insert = ['\n'] # Ensure at least one newline for empty text

            new_file_content = lines[:start_idx] + new_lines_to_insert + lines[end_idx + 1:]

            with open(full_path, 'w') as f:
                f.writelines(new_file_content)
            return ""
        except FileNotFoundError:
            return f"File '{file_name}' not found for editing."
        except Exception as e:
            return f"An error occurred while editing file: {e}"

    def edit_all(self, file_name: str, text: str) -> str:
        """
        Replaces the entire content of a file with the provided text.
        """
        full_path = self._get_full_path(file_name)
        try:
            with open(full_path, 'w') as f:
                f.write(text)
            return ""
        except Exception as e:
            return f"An error occurred while editing all of file: {e}"

    # --- Create new file operations ---
    def create_new_file_only(self, file_name: str) -> str:
        """
        Creates a new empty file.
        """
        full_path = self._get_full_path(file_name)
        if os.path.exists(full_path):
            return f"File '{file_name}' already exists."
        try:
            with open(full_path, 'x') as f:
                pass
            return ""
        except FileExistsError: # Redundant due to above check, but good for robustness
            return f"File '{file_name}' already exists unexpectedly."
        except Exception as e:
            return f"An error occurred while creating new file: {e}"

    def create_new_file_and_text(self, file_name: str, text: str) -> str:
        """
        Creates a new file and writes the provided text to it.
        Ensures the directory path exists before creating the file.
        """
        full_path = self._get_full_path(file_name)
        
        # Extract the directory path from the full file path
        directory = os.path.dirname(full_path)
        
        try:
            # Create the directory if it doesn't exist (and any parent directories)
            # exist_ok=True prevents an error if the directory already exists
            if directory and not os.path.exists(directory):
                os.makedirs(directory, exist_ok=True) # This is the key line for cascade creation

            with open(full_path, 'w') as f:
                f.write(text)
            return ""
        except Exception as e:
            return f"An error occurred while creating new file with text: {e}"

    def delete_file(self, file_name: str) -> str:
        """
        Deletes a file from the managed directory.
        """
        full_path = self._get_full_path(file_name)
        try:
            if os.path.exists(full_path):
                os.remove(full_path)
                return ""
            else:
                return f"File '{file_name}' not found."
        except Exception as e:
            return f"An error occurred while deleting file: {e}"

    def list_files(self) -> tuple[list, str]:
        """
        Lists all files in the managed directory.
        """
        try:
            files = [f for f in os.listdir(self.base_dir) if os.path.isfile(self._get_full_path(f))]
            return files, ""
        except Exception as e:
            return [], f"An error occurred while listing files: {e}"

    # --- Folder operations ---
    def create_folder(self, folder_name: str) -> str:
        """
        Creates a new folder within the managed directory.
        """
        full_path = os.path.join(self.base_dir, folder_name)
        try:
            if os.path.exists(full_path):
                return f"Folder '{folder_name}' already exists."
            os.makedirs(full_path)
            return ""
        except Exception as e:
            return f"An error occurred while creating folder: {e}"
